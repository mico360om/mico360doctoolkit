"""AI metadata suggestions: read a document, propose document properties.

The document's own text is extracted locally (PyMuPDF / python-docx / pptx) and
only an excerpt is sent to the configured AI provider — never the whole file.
The model is asked for strict JSON, and the reply is parsed defensively because
small models like to wrap JSON in prose or code fences.

Nothing is ever written to the file here: this module only *suggests*. Applying
is the user's explicit choice in the UI.
"""
from __future__ import annotations

import json
import re
from pathlib import Path

from mico360.core.ai import AiConfig, AiError, chat
from mico360.logging_setup import get_logger

log = get_logger("mico360.ai_metadata")

# Every metadata field the AI may propose, in the order the UI shows them.
# These keys match the Edit Metadata tool's option keys exactly, so a suggestion
# can be written straight into the matching control.
FIELDS = ("title", "author", "subject", "keywords", "creator", "producer",
          "creation_date", "mod_date", "company", "manager", "category",
          "comments", "custom", "copyright", "language", "trapped")

FIELD_LABELS = {
    "title": "Title",
    "author": "Author",
    "subject": "Subject",
    "keywords": "Keywords",
    "creator": "Creator (authoring app)",
    "producer": "Producer (PDF software)",
    "creation_date": "Creation date",
    "mod_date": "Modification date",
    "company": "Company",
    "manager": "Manager",
    "category": "Category",
    "comments": "Description / Comments",
    "custom": "Custom properties",
    "copyright": "Copyright",
    "language": "Language",
    "trapped": "Trapped",
}

# Trapped is a fixed choice, not free text.
_TRAPPED_VALUES = {"true": "True", "false": "False", "unknown": "Unknown"}

# How much document text to send. Enough to characterise the document, well
# under the platform's 32,000-character prompt limit.
MAX_CHARS = 8000

_PROMPT = """You are a document cataloguing assistant. Read the document excerpt and return metadata for it.

Reply with ONE JSON object and nothing else. Use exactly these keys:
"title", "author", "subject", "keywords", "creator", "producer",
"creation_date", "mod_date", "company", "manager", "category", "comments",
"custom", "copyright", "language", "trapped".

Rules for each field:
- title: the document's real title, not the file name. Concise.
- author: the person or team who wrote it, if the document says so.
- subject: one short line describing what the document is about.
- keywords: 3-8 comma-separated terms, lowercase.
- creator: the application the document was authored in, only if it is evident.
- producer: the software that produced the PDF, only if it is evident.
- creation_date: the date PRINTED IN the document (e.g. an invoice or report
  date) as YYYY-MM-DD. Never guess, and never use today's date.
- mod_date: a revision date printed in the document as YYYY-MM-DD, else "".
- company: the organisation the document belongs to, if named.
- manager: the responsible manager or owner, if named.
- category: one short classifier, e.g. Report, Invoice, Contract, Manual.
- comments: a 1-2 sentence description.
- custom: extra properties worth recording, one per line as "Key = Value"
  (e.g. "Invoice Number = INV-1042"). Use "" if there are none.
- copyright: a copyright line if the document states one.
- language: the BCP-47 tag of the document's language, e.g. en-US, ar, fr-FR.
- trapped: "True", "False" or "Unknown" — use "" unless the document says.

CRITICAL: use "" for anything you cannot determine from the excerpt. Never
invent a value, a person or a date. An empty string is always better than a
guess, because a blank suggestion leaves the existing metadata untouched.

Document excerpt:
---
{excerpt}
---"""


class NoTextError(AiError):
    """The document has no extractable text (e.g. a scan needing OCR first)."""


# =====================================================================
# Local text extraction
# =====================================================================
def extract_text(path: Path, max_chars: int = MAX_CHARS) -> str:
    """Pull a representative text excerpt out of a document, locally.

    Raises NoTextError when there's nothing to read — a scanned PDF should be
    run through Searchable PDF (OCR) first, and we say so.
    """
    p = Path(path)
    ext = p.suffix.lower()
    try:
        if ext == ".pdf":
            text = _from_pdf(p, max_chars)
        elif ext == ".docx":
            import docx
            text = "\n".join(par.text for par in docx.Document(str(p)).paragraphs)
        elif ext == ".pptx":
            text = _from_pptx(p)
        elif ext in (".txt", ".md", ".csv"):
            text = p.read_text(encoding="utf-8", errors="replace")
        else:
            raise NoTextError(
                f"'{p.name}' isn't a text document the AI can read "
                "(PDF, Word, PowerPoint, text).")
    except NoTextError:
        raise
    except Exception as exc:      # noqa: BLE001
        raise AiError(f"Couldn't read '{p.name}' ({exc}).")

    text = re.sub(r"[ \t]+", " ", text or "")
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    if len(text) < 20:
        if ext == ".pdf":
            raise NoTextError(
                f"'{p.name}' has no selectable text — it looks like a scan. Run "
                "Searchable PDF (OCR) on it first, then try again.")
        raise NoTextError(
            f"'{p.name}' has no readable text — its content appears to be "
            "images. Convert it with Office → PDF, run Searchable PDF "
            "(OCR) on the result, then try again.")
    return text[:max_chars]


def _from_pptx(p: Path) -> str:
    """Slide text, including grouped shapes, tables and speaker notes — a real
    deck often puts its wording in those rather than plain text boxes."""
    from pptx import Presentation

    chunks: list[str] = []

    def walk(shapes):
        for shape in shapes:
            try:
                if str(getattr(shape, "shape_type", "")).startswith("GROUP"):
                    walk(shape.shapes)
                    continue
                if getattr(shape, "has_text_frame", False):
                    chunks.append(shape.text_frame.text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        chunks.append(" ".join(c.text for c in row.cells))
            except Exception:          # noqa: BLE001 - one odd shape mustn't stop us
                continue

    for slide in Presentation(str(p)).slides:
        walk(slide.shapes)
        try:
            if slide.has_notes_slide:
                chunks.append(slide.notes_slide.notes_text_frame.text)
        except Exception:              # noqa: BLE001
            pass
    return "\n".join(c for c in chunks if c)


def _from_pdf(p: Path, max_chars: int) -> str:
    """First pages carry the identity of a document; that's what we need."""
    import fitz
    doc = fitz.open(str(p))
    try:
        parts = []
        total = 0
        for i in range(min(doc.page_count, 10)):
            t = doc[i].get_text("text")
            if t:
                parts.append(t)
                total += len(t)
            if total >= max_chars:
                break
        return "\n".join(parts)
    finally:
        doc.close()


# =====================================================================
# Suggestion
# =====================================================================
def _parse_json_object(reply: str) -> dict:
    """Extract the JSON object from a model reply that may include prose or
    ```json fences."""
    s = (reply or "").strip()
    fence = re.search(r"```(?:json)?\s*(.+?)```", s, re.S)
    if fence:
        s = fence.group(1).strip()
    if not s.startswith("{"):
        start, end = s.find("{"), s.rfind("}")
        if start >= 0 and end > start:
            s = s[start:end + 1]
    try:
        obj = json.loads(s)
    except (ValueError, TypeError):
        raise AiError("The AI reply couldn't be understood as metadata. "
                      "Try again, or use a larger model.")
    return obj if isinstance(obj, dict) else {}


_PLACEHOLDERS = ("n/a", "none", "null", "not specified", "not stated",
                 "not available", "unspecified", "-", "--", "tbd", "todo",
                 "no author", "no date")


def _clean(value, field: str = "") -> str:
    """Normalise one suggested value. Returns "" for anything unusable, which
    means the field is left as-is rather than being blanked."""
    if value is None:
        return ""
    if isinstance(value, dict):                    # custom props as an object
        value = "\n".join(f"{k} = {v}" for k, v in value.items())
    if isinstance(value, (list, tuple)):
        sep = "\n" if field == "custom" else ", "
        value = sep.join(str(v) for v in value)
    text = str(value)

    if field == "custom":
        # Keep one "Key = Value" per line; drop anything that isn't a pair.
        lines = [re.sub(r"\s+", " ", ln).strip() for ln in text.splitlines()]
        keep = [ln for ln in lines
                if "=" in ln and ln.split("=", 1)[0].strip()
                and ln.split("=", 1)[1].strip()
                and ln.split("=", 1)[1].strip().lower() not in _PLACEHOLDERS]
        return "\n".join(keep)[:1000]

    text = re.sub(r"\s+", " ", text).strip()
    low = text.lower().strip(" .")
    if not text or low in _PLACEHOLDERS:
        return ""

    if field == "trapped":
        return _TRAPPED_VALUES.get(low, "")        # only the three valid values
    if field in ("creation_date", "mod_date"):
        m = re.search(r"(\d{4})[-/](\d{1,2})[-/](\d{1,2})", text)
        if not m:
            return ""                              # unusable date -> keep existing
        y, mo, d = (int(x) for x in m.groups())
        if not (1900 <= y <= 2200 and 1 <= mo <= 12 and 1 <= d <= 31):
            return ""
        return f"{y:04d}-{mo:02d}-{d:02d}"
    if field == "language":
        return text if re.fullmatch(r"[A-Za-z]{2,3}(-[A-Za-z0-9]{2,8})*", text) else ""
    return text[:500]


def suggest_metadata(path: Path, cfg: AiConfig, excerpt: str | None = None,
                     cancel=None) -> dict:
    """Return {field: suggested value} for a document. Raises AiError.

    ``cancel`` is an optional event-like object (``is_set()``/``wait()``) passed
    through to the request layer so a user's Cancel interrupts a retry wait.
    """
    text = excerpt if excerpt is not None else extract_text(Path(path))
    reply = chat(cfg, [
        {"role": "system",
         "content": "You reply with a single JSON object and no other text."},
        {"role": "user", "content": _PROMPT.format(excerpt=text)},
    ], cancel=cancel)
    obj = _parse_json_object(reply)
    out = {}
    for key in FIELDS:
        val = _clean(obj.get(key), key)
        if val:
            out[key] = val
    if not out:
        raise AiError("The AI didn't return any usable metadata for this "
                      "document. Try again, or use a larger model.")
    log.info("AI suggested %d metadata field(s) for %s", len(out), Path(path).name)
    return out
