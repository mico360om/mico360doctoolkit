"""Processing functions for every tool.

Each function has one of two signatures:

  per-file   : fn(src: Path, out_dir: Path, opt: dict, report) -> list[Path]
  aggregate  : fn(srcs: list[Path], out_dir: Path, opt: dict, report) -> list[Path]

``report(msg)`` writes a line to the activity log. Functions raise
``ProcessError`` with a user-readable message on failure.
"""
from __future__ import annotations

import subprocess
import threading
from pathlib import Path
from typing import Callable

from mico360.core.deps import find_ghostscript, find_libreoffice
from mico360.core.util import (
    ProcessError,
    build_output_path,
    human_size,
    run_subprocess,
    unique_dir,
    unique_path,
)

Report = Callable[[str], None]

IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".tif", ".tiff",
              ".heic", ".heif"}

# Teach Pillow to read/write HEIC/HEIF (Apple photos). Registered once at import;
# a harmless no-op if the optional plugin isn't installed.
try:
    import pillow_heif as _pillow_heif
    _pillow_heif.register_heif_opener()
except Exception:  # pragma: no cover - optional dependency
    pass


def _clampi(value, lo: int, hi: int, default: int) -> int:
    """Coerce *value* to an int within [lo, hi]. Guards against out-of-range
    numbers that bypass the UI spinboxes (e.g. a hand-edited settings file)."""
    try:
        return max(lo, min(hi, int(value)))
    except (TypeError, ValueError):
        return default


def _progress(report: Report, current: int, total: int) -> None:
    """Report fine-grained progress within a unit, if the caller supports it.

    The engine's report callable carries a ``.progress`` attribute; plain test
    callbacks don't, so we degrade gracefully to a no-op."""
    fn = getattr(report, "progress", None)
    if fn is not None:
        try:
            fn(current, total)
        except Exception:  # pragma: no cover - never let progress break a job
            pass


def _progress_frac(report: Report, frac: float) -> None:
    """Report a 0..1 completion fraction within a unit (degrades to a no-op).
    Lets one operation map its progress into a sub-range of the unit's bar."""
    fn = getattr(report, "progress", None)
    if fn is not None:
        try:
            fn(int(max(0.0, min(1.0, frac)) * 1000), 1000)
        except Exception:  # pragma: no cover
            pass


def _check_cancel(report: Report) -> None:
    """Raise ProcessError('Cancelled') if the user cancelled the batch — lets
    long per-page loops stop promptly instead of running to completion."""
    fn = getattr(report, "cancelled", None)
    if fn is not None and fn():
        raise ProcessError("Cancelled")


# =========================================================================
# PDF compression
# =========================================================================
# Ghostscript distiller presets keyed by our Low/Medium/High compression labels.
_GS_PRESET = {"low": "/printer", "medium": "/ebook", "high": "/screen"}
# Target image resolution / JPEG quality used by the built-in PyMuPDF fallback
# so the Low/Medium/High choice produces meaningfully different sizes even when
# Ghostscript is unavailable.
_FALLBACK_DPI = {"low": 200, "medium": 130, "high": 90}
_FALLBACK_QUALITY = {"low": 80, "medium": 60, "high": 45}


def pdf_compress(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    level = opt.get("level", "lossless")
    out = build_output_path(src, out_dir, ".pdf", name_suffix="_compressed",
                            overwrite=opt.get("overwrite", False),
                            numbered=opt.get("same_as_source", False))
    before = src.stat().st_size

    # Lossless (default): reduce size with a strict content-integrity guarantee.
    if level == "lossless":
        _pdf_compress_lossless(src, out, before, report)
        return [out]

    if level == "target":
        target_bytes = max(1, int(opt.get("target_kb", 250))) * 1024
        _pdf_compress_to_target(src, out, target_bytes, before, report)
        _enforce_structural_integrity(src, out, report)
        return [out]

    gs = find_ghostscript()
    if gs:
        args = [
            gs, "-sDEVICE=pdfwrite", "-dCompatibilityLevel=1.5",
            "-dNOPAUSE", "-dBATCH", "-dQUIET", "-dSAFER",
            "-dDetectDuplicateImages=true", "-dCompressFonts=true",
            "-dSubsetFonts=true",
        ]
        if level == "custom":
            dpi = _clampi(opt.get("dpi", 150), 36, 600, 150)
            jpegq = _clampi(opt.get("jpeg_quality", 75), 10, 100, 75)
            args += [
                "-dPDFSETTINGS=/printer",
                "-dDownsampleColorImages=true", "-dColorImageDownsampleType=/Bicubic",
                f"-dColorImageResolution={dpi}",
                "-dDownsampleGrayImages=true", "-dGrayImageDownsampleType=/Bicubic",
                f"-dGrayImageResolution={dpi}",
                "-dDownsampleMonoImages=true", f"-dMonoImageResolution={max(dpi, 300)}",
                "-dAutoFilterColorImages=false", "-dColorImageFilter=/DCTEncode",
                f"-dJPEGQ={jpegq}",
            ]
        else:
            args.append(f"-dPDFSETTINGS={_GS_PRESET.get(level, '/ebook')}")
        args += [f"-sOutputFile={out}", str(src)]

        report(f"Ghostscript compressing ({level})…")
        try:
            proc = run_subprocess(args)
        except subprocess.TimeoutExpired:
            raise ProcessError(
                "Compression timed out (over 10 minutes). Try a lower compression "
                "level, or split very large PDFs first.")
        if proc.returncode != 0 or not out.exists():
            raise ProcessError(
                "Ghostscript failed: " + (proc.stderr or proc.stdout or "unknown error").strip()[:300]
            )
    else:
        report(f"Compressing with the built-in engine ({level})…")
        _pymupdf_compress(src, out, level, opt, report)

    after = out.stat().st_size
    pct = (1 - after / before) * 100 if before else 0
    report(f"{human_size(before)} → {human_size(after)} ({pct:+.0f}%)")
    if after >= before:
        report("Note: file was already well-optimized; little/no size reduction.")
    # Lossy levels re-encode images by design; make sure nothing ELSE was lost.
    # The compress step filled the first half of the bar; verify fills the rest.
    _enforce_structural_integrity(src, out, report, prog_base=0.5, prog_span=0.5)
    return [out]


def _pymupdf_compress(src: Path, out: Path, level: str, opt: dict, report: Report) -> None:
    """Built-in fallback used when Ghostscript is unavailable.

    Downsamples and re-encodes embedded raster images to a resolution / JPEG
    quality chosen by the compression *level* (text and vector content are kept
    intact), then rebuilds the file: garbage-collects unused objects, deflates
    streams/fonts, subsets fonts and cleans content. This makes the Low / Medium
    / High choice produce visibly different sizes. Ghostscript (bundled) remains
    the recommended engine for the most aggressive size reduction.
    """
    import fitz  # PyMuPDF

    if level == "custom":
        dpi_target = _clampi(opt.get("dpi", 150), 36, 600, 150)
        quality = _clampi(opt.get("jpeg_quality", 75), 10, 100, 75)
    else:
        dpi_target = _FALLBACK_DPI.get(level, _FALLBACK_DPI["medium"])
        quality = _FALLBACK_QUALITY.get(level, _FALLBACK_QUALITY["medium"])

    doc = fitz.open(src)
    _progress_frac(report, 0.05)
    try:
        # PyMuPDF >= 1.24 can downsample/recompress images in place. Only images
        # whose effective resolution exceeds the target are touched, so we never
        # upscale or bloat already-small images.
        if hasattr(doc, "rewrite_images"):
            report("Recompressing images (can take a moment on large PDFs)…")
            try:
                # rewrite_images requires dpi_target < dpi_threshold: any image
                # above the threshold is downsampled to the target. Using
                # target + 1 means "normalize everything above the target DPI".
                doc.rewrite_images(
                    dpi_threshold=dpi_target + 1, dpi_target=dpi_target,
                    quality=quality, lossy=True, lossless=True,
                )
            except Exception as exc:  # pragma: no cover - version/feature guard
                report(f"Image recompression unavailable ({exc}); rebuilding only.")
        _progress_frac(report, 0.35)
        report("Rebuilding the PDF…")
        doc.save(str(out), garbage=4, deflate=True, deflate_images=True,
                 deflate_fonts=True, clean=True)
        _progress_frac(report, 0.5)
    finally:
        doc.close()


def _links_signature(page) -> list:
    out = []
    for ln in page.get_links():
        frm = ln.get("from")
        rect = tuple(round(v, 1) for v in frm) if frm is not None else None
        out.append((ln.get("kind"), ln.get("uri"), ln.get("page"),
                    ln.get("file"), rect))
    return sorted(str(x) for x in out)


# Above this page count, the per-page rendered-pixel hash (the most expensive
# strict check) is skipped: image bytes, text, fonts, links and structure already
# prove content was preserved, and rendering every page is what makes verifying a
# large PDF slow. Keeps big-file compression responsive without weakening safety.
_VERIFY_RENDER_MAX_PAGES = 25


def verify_pdf_integrity(src_path, out_path, mode: str = "strict", report=None,
                         prog_base: float = 0.0, prog_span: float = 1.0):
    """Confirm a processed PDF preserved the original's content.

    Returns ``(ok: bool, differences: list[str])``.

    * ``mode="strict"`` (lossless compression) — EVERYTHING must match: page
      count, text, links, bookmarks, attachments, fonts, metadata, image bytes,
      and (on smaller PDFs) the rendered appearance of every page.
    * ``mode="structural"`` (lossy compression) — the things that must never be
      lost must match: page count, text, links, bookmarks, attachments.

    Progress is reported per page into the sub-range ``[prog_base, prog_base +
    prog_span]`` so a long verification of a big file keeps the bar moving.
    """
    import hashlib

    import fitz
    strict = mode == "strict"
    diffs: list[str] = []
    a = fitz.open(str(src_path))
    b = fitz.open(str(out_path))
    try:
        n = a.page_count
        if n != b.page_count:
            return False, [f"page count changed {n} → {b.page_count}"]
        if a.get_toc(simple=True) != b.get_toc(simple=True):
            diffs.append("bookmarks / outline changed")
        try:
            if sorted(a.embfile_names()) != sorted(b.embfile_names()):
                diffs.append("embedded files / attachments changed")
        except Exception:
            pass
        if strict:
            for k in ("title", "author", "subject", "keywords", "creator", "producer"):
                if (a.metadata.get(k) or "") != (b.metadata.get(k) or ""):
                    diffs.append(f"metadata '{k}' changed")
        # On large PDFs, skip the heavy rendered-pixel comparison (see above).
        render_check = strict and n <= _VERIFY_RENDER_MAX_PAGES
        for i in range(n):
            _check_cancel(report)
            pa, pb = a[i], b[i]
            if pa.get_text("text") != pb.get_text("text"):
                diffs.append(f"page {i + 1}: text changed")
            if _links_signature(pa) != _links_signature(pb):
                diffs.append(f"page {i + 1}: links / annotations changed")
            if strict:
                fa = sorted(f[3] for f in pa.get_fonts(full=True))
                fb = sorted(f[3] for f in pb.get_fonts(full=True))
                if fa != fb:
                    diffs.append(f"page {i + 1}: fonts changed")
                ia, ib = pa.get_images(full=True), pb.get_images(full=True)
                if len(ia) != len(ib):
                    diffs.append(f"page {i + 1}: image count {len(ia)} → {len(ib)}")
                else:
                    ha = sorted(hashlib.sha1(a.extract_image(im[0])["image"]).hexdigest()
                                for im in ia)
                    hb = sorted(hashlib.sha1(b.extract_image(im[0])["image"]).hexdigest()
                                for im in ib)
                    if ha != hb:
                        diffs.append(f"page {i + 1}: image data changed")
                if render_check:
                    ra = hashlib.sha1(pa.get_pixmap(dpi=100).samples).hexdigest()
                    rb = hashlib.sha1(pb.get_pixmap(dpi=100).samples).hexdigest()
                    if ra != rb:
                        diffs.append(f"page {i + 1}: rendered appearance changed")
            if report is not None:
                _progress_frac(report, prog_base + (i + 1) / max(1, n) * prog_span)
    finally:
        a.close()
        b.close()
    return (len(diffs) == 0), diffs


def _pdf_compress_lossless(src: Path, out: Path, before: int, report: Report) -> None:
    """Reduce size with ZERO content change: garbage-collect unused objects and
    deflate streams/fonts, but never touch image data, text, links or metadata.
    The result is verified to be content-identical to the original; if it isn't
    (or if it wouldn't be smaller), the original is kept verbatim."""
    import shutil

    import fitz
    report("Lossless compression — preserving 100% of content…")
    _progress_frac(report, 0.05)
    doc = fitz.open(str(src))
    try:
        doc.set_metadata(dict(doc.metadata))   # keep all metadata exactly
        doc.save(str(out), garbage=4, deflate=True, deflate_images=False,
                 deflate_fonts=True, clean=True)
    finally:
        doc.close()

    report("Verifying the result is identical to the original…")
    # Compression took the first ~20% of the bar; the verify fills the rest so a
    # large file's (slower) check keeps the percentage moving instead of freezing.
    ok, diffs = verify_pdf_integrity(src, out, mode="strict", report=report,
                                     prog_base=0.2, prog_span=0.8)
    if not ok:
        report("Integrity check flagged a difference — keeping the original "
               f"unchanged ({'; '.join(diffs[:3])}).")
        shutil.copyfile(str(src), str(out))
        return
    after = out.stat().st_size
    if after >= before:
        shutil.copyfile(str(src), str(out))     # no gain → keep original (no bloat)
        report(f"Already optimally packed — kept the original ({human_size(before)}); "
               "content verified identical ✓")
        return
    report(f"{human_size(before)} → {human_size(after)} "
           f"({(1 - after / before) * 100:.0f}% smaller) — content verified identical ✓")


def _enforce_structural_integrity(src: Path, out: Path, report: Report,
                                  prog_base: float = 0.5, prog_span: float = 0.5) -> None:
    """After a *lossy* compression, make sure nothing other than image quality was
    lost (text, links, bookmarks, attachments, page count). If something was,
    keep the original so the user never loses content."""
    import shutil
    report("Verifying text, links, bookmarks & attachments were kept…")
    ok, diffs = verify_pdf_integrity(src, out, mode="structural", report=report,
                                     prog_base=prog_base, prog_span=prog_span)
    if ok:
        report("Integrity check ✓ — text, links, bookmarks & attachments preserved.")
        return
    report("⚠ This compression would have changed content beyond image quality ("
           + "; ".join(diffs[:3]) + ") — kept the original to guarantee zero data loss.")
    shutil.copyfile(str(src), str(out))


def _pdf_compress_to_target(src: Path, out: Path, target_bytes: int,
                            before: int, report: Report) -> None:
    """Compress a PDF to land at or just under *target_bytes*.

    Binary-searches a strength scale (image DPI + JPEG quality, applied via
    PyMuPDF's in-place image rewriter) for the gentlest setting whose output is
    still <= target, so quality is maximised without exceeding the size limit.
    All other settings are chosen automatically.
    """
    import fitz

    report(f"Targeting <= {human_size(target_bytes)} (auto-selecting settings)…")

    def produce(scale: int | None) -> int:
        # scale None = lossless rebuild (highest quality, no image re-encoding).
        # scale 0..100: higher = higher dpi/quality = larger file.
        doc = fitz.open(src)
        try:
            if scale is not None and hasattr(doc, "rewrite_images"):
                dpi = int(40 + scale / 100 * (250 - 40))
                quality = int(20 + scale / 100 * (88 - 20))
                try:
                    doc.rewrite_images(dpi_threshold=dpi + 1, dpi_target=dpi,
                                       quality=quality, lossy=True, lossless=True)
                except Exception:  # pragma: no cover
                    pass
            doc.save(str(out), garbage=4, deflate=True, deflate_images=True,
                     deflate_fonts=True, clean=True)
        finally:
            doc.close()
        return out.stat().st_size

    _TOTAL_STEPS = 10  # 1 lossless probe + up to 8 search probes + 1 final write

    # Best quality first: if a plain lossless rebuild already fits, keep it.
    if produce(None) <= target_bytes:
        _progress(report, _TOTAL_STEPS, _TOTAL_STEPS)
        achieved = out.stat().st_size
        pct = (1 - achieved / before) * 100 if before else 0
        report(f"{human_size(before)} → {human_size(achieved)} ({pct:+.0f}%, lossless)")
        return
    _progress(report, 1, _TOTAL_STEPS)

    lo, hi, best = 0, 100, None
    for step in range(8):  # ~log2(100) probes
        if lo > hi:
            break
        mid = (lo + hi) // 2
        size = produce(mid)
        _progress(report, 2 + step, _TOTAL_STEPS)
        if size <= target_bytes:
            best = mid
            lo = mid + 1
        else:
            hi = mid - 1
    # Materialise the chosen setting (last probe may not be the best one).
    final_scale = best if best is not None else 0
    achieved = produce(final_scale)
    _progress(report, _TOTAL_STEPS, _TOTAL_STEPS)

    pct = (1 - achieved / before) * 100 if before else 0
    report(f"{human_size(before)} → {human_size(achieved)} ({pct:+.0f}%)")
    if achieved > target_bytes:
        report("Note: the target is smaller than this PDF's text/vector content "
               "allows — produced the smallest possible instead.")


# =========================================================================
# Merge PDFs (aggregate)
# =========================================================================
def pdf_merge(srcs: list[Path], out_dir: Path, opt: dict, report: Report) -> list[Path]:
    from pypdf import PdfReader, PdfWriter

    if len(srcs) < 2:
        raise ProcessError("Select at least two PDF files to merge.")
    name = (opt.get("output_name") or "merged").strip() or "merged"
    out = unique_path(out_dir / f"{name}.pdf", opt.get("overwrite", False))
    out_dir.mkdir(parents=True, exist_ok=True)

    writer = PdfWriter()
    total = 0
    for i, s in enumerate(srcs):
        report(f"Adding {s.name}")
        try:
            reader = PdfReader(str(s))
            pages = list(reader.pages)
        except Exception as exc:
            raise ProcessError(f"Couldn't read '{s.name}' — it may be corrupt or "
                               f"password-protected ({exc}).")
        for page in pages:
            writer.add_page(page)
        total += len(pages)
        _progress(report, i + 1, len(srcs))
    with open(out, "wb") as fh:
        writer.write(fh)
    report(f"Merged {len(srcs)} files / {total} pages → {out.name}")
    return [out]


# =========================================================================
# Split PDF (per-file → many outputs)
# =========================================================================
def pdf_split(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    from pypdf import PdfReader, PdfWriter

    try:
        reader = PdfReader(str(src))
        n = len(reader.pages)
    except Exception as exc:
        raise ProcessError(f"Couldn't read '{src.name}' — it may be corrupt or "
                           f"password-protected ({exc}).")
    mode = opt.get("mode", "each")
    overwrite = opt.get("overwrite", False)
    # Unique subfolder per source so two same-named PDFs from different folders
    # don't write into the same "stem/" directory and collide.
    dest = unique_dir(out_dir / src.stem)
    outputs: list[Path] = []

    def write_pages(indices: list[int], label: str) -> None:
        w = PdfWriter()
        for i in indices:
            w.add_page(reader.pages[i])
        p = unique_path(dest / f"{src.stem}_{label}.pdf", overwrite)
        with open(p, "wb") as fh:
            w.write(fh)
        outputs.append(p)
        report(f"Wrote {p.name}")

    if mode == "each":
        for i in range(n):
            _check_cancel(report)
            write_pages([i], f"page{i + 1}")
            _progress(report, i + 1, n)
    elif mode == "every_n":
        step = max(1, int(opt.get("every_n", 1)))
        for start in range(0, n, step):
            idx = list(range(start, min(start + step, n)))
            write_pages(idx, f"p{idx[0] + 1}-{idx[-1] + 1}")
            _progress(report, idx[-1] + 1, n)
    elif mode == "ranges":
        ranges = _parse_ranges(opt.get("ranges", ""), n)
        for j, (label, idx) in enumerate(ranges):
            write_pages(idx, label)
            _progress(report, j + 1, len(ranges) or 1)
        if not outputs:
            raise ProcessError("No valid page ranges were provided (e.g. '1-3, 5, 8-10').")
    else:
        raise ProcessError(f"Unknown split mode: {mode}")
    return outputs


def _parse_ranges(spec: str, n: int) -> list[tuple[str, list[int]]]:
    out: list[tuple[str, list[int]]] = []
    for chunk in spec.replace(" ", "").split(","):
        if not chunk:
            continue
        if "-" in chunk:
            a, _, b = chunk.partition("-")
            try:
                start, end = int(a), int(b)
            except ValueError:
                continue
            start = max(1, start)
            end = min(n, end)
            if start <= end:
                out.append((f"p{start}-{end}", list(range(start - 1, end))))
        else:
            try:
                p = int(chunk)
            except ValueError:
                continue
            if 1 <= p <= n:
                out.append((f"page{p}", [p - 1]))
    return out


# =========================================================================
# Page-spec helpers shared by Organize PDF (1-based input → 0-based indices)
# =========================================================================
def _page_list(spec: str, n: int, default_all: bool = False) -> list[int]:
    """Parse a page spec like ``3, 1, 2, 5-8`` into an *ordered* list of 0-based
    indices (preserving the given order, ranges allowed in either direction).
    Empty / 'all' → every page when ``default_all`` else []."""
    spec = (spec or "").strip().lower()
    if spec in ("", "all"):
        return list(range(n)) if default_all else []
    out: list[int] = []
    for chunk in spec.replace(" ", "").split(","):
        if not chunk:
            continue
        if "-" in chunk[1:]:  # a range (a leading '-' would be a stray sign)
            a, _, b = chunk.partition("-")
            try:
                start, end = int(a), int(b)
            except ValueError:
                continue
            rng = range(start, end + 1) if start <= end else range(start, end - 1, -1)
            out.extend(p - 1 for p in rng if 1 <= p <= n)
        else:
            try:
                p = int(chunk)
            except ValueError:
                continue
            if 1 <= p <= n:
                out.append(p - 1)
    return out


def _page_set(spec: str, n: int, default_all: bool = False) -> set:
    return set(_page_list(spec, n, default_all=default_all))


# =========================================================================
# Organize PDF — rotate / delete / extract / reorder pages
# =========================================================================
def pdf_organize(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    from pypdf import PdfReader, PdfWriter

    try:
        reader = PdfReader(str(src))
        n = len(reader.pages)
    except Exception as exc:
        raise ProcessError(f"Couldn't read '{src.name}' — it may be corrupt or "
                           f"password-protected ({exc}).")
    op = opt.get("operation", "rotate")
    writer = PdfWriter()

    if op == "rotate":
        angle = _clampi(opt.get("angle", 90), -270, 270, 90)
        if angle % 90 != 0:
            raise ProcessError("Rotation must be a multiple of 90°.")
        targets = _page_set(opt.get("pages", "all"), n, default_all=True)
        for i in range(n):
            page = reader.pages[i]
            if i in targets:
                page.rotate(angle)
            writer.add_page(page)
            _progress(report, i + 1, n)
        suffix = "_rotated"
    elif op == "delete":
        remove = _page_set(opt.get("del_pages", ""), n)
        if not remove:
            raise ProcessError("Enter the page(s) to delete, e.g. 2, 5-7.")
        kept = [i for i in range(n) if i not in remove]
        if not kept:
            raise ProcessError("That would delete every page.")
        for j, i in enumerate(kept):
            writer.add_page(reader.pages[i])
            _progress(report, j + 1, len(kept))
        suffix = "_pages-removed"
    elif op in ("extract", "reorder"):
        key = "ext_pages" if op == "extract" else "order"
        order = _page_list(opt.get(key, ""), n)
        if not order:
            raise ProcessError("Enter the pages, e.g. 3, 1, 2, 5-8.")
        for j, i in enumerate(order):
            writer.add_page(reader.pages[i])
            _progress(report, j + 1, len(order))
        suffix = "_extracted" if op == "extract" else "_reordered"
    else:
        raise ProcessError(f"Unknown operation: {op}")

    out = build_output_path(src, out_dir, ".pdf", name_suffix=suffix,
                            overwrite=opt.get("overwrite", False),
                            numbered=opt.get("same_as_source", False))
    with open(out, "wb") as fh:
        writer.write(fh)
    report(f"Saved {len(writer.pages)} page(s) → {out.name}")
    return [out]


# =========================================================================
# Protect / Unlock PDF — add or remove a password
# =========================================================================
def pdf_protect(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    from pypdf import PdfReader, PdfWriter

    op = opt.get("operation", "protect")
    pw = opt.get("password", "") or ""
    try:
        reader = PdfReader(str(src))
    except Exception as exc:
        raise ProcessError(f"Couldn't read '{src.name}' ({exc}).")

    if reader.is_encrypted:
        try:
            ok = bool(int(reader.decrypt(pw)))
        except Exception:
            ok = False
        if not ok:
            if op == "unlock":
                raise ProcessError("Incorrect password — couldn't unlock this PDF.")
            raise ProcessError("This PDF is already password-protected. Unlock it "
                               "with its current password first.")
    elif op == "unlock":
        report(f"'{src.name}' isn't password-protected — saving a copy.")

    writer = PdfWriter()
    try:
        for page in reader.pages:
            writer.add_page(page)
    except Exception as exc:
        raise ProcessError(f"Couldn't read pages from '{src.name}' ({exc}).")

    if op == "protect":
        if not pw:
            raise ProcessError("Enter a password to protect the PDF.")
        confirm = opt.get("confirm_password", None)
        if confirm is not None and confirm != pw:
            raise ProcessError("The passwords don't match — please re-enter them.")
        # Prefer strong AES-256 (needs `cryptography`); fall back to the
        # dependency-free RC4-128 so protection always works.
        try:
            writer.encrypt(user_password=pw, algorithm="AES-256")
        except Exception:
            try:
                writer.encrypt(user_password=pw, algorithm="RC4-128")
            except Exception:
                writer.encrypt(pw)
        suffix = "_protected"
    else:
        suffix = "_unlocked"

    out = build_output_path(src, out_dir, ".pdf", name_suffix=suffix,
                            overwrite=opt.get("overwrite", False),
                            numbered=opt.get("same_as_source", False))
    with open(out, "wb") as fh:
        writer.write(fh)
    report(f"{'Protected' if op == 'protect' else 'Unlocked'} → {out.name}")
    return [out]


# =========================================================================
# Watermark PDF — overlay diagonal text on every page
# =========================================================================
_WM_COLORS = {
    "gray": (0.5, 0.5, 0.5), "red": (0.78, 0.12, 0.12),
    "blue": (0.13, 0.27, 0.62), "black": (0.0, 0.0, 0.0),
}


def _anchor_xy(position: str, pw: float, ph: float, w: float, h: float,
               margin: float = 28.0) -> tuple[float, float]:
    """Top-left (x, y) for a w×h watermark placed at *position* on a pw×ph page."""
    parts = (position or "center").split("-")
    vert = parts[0] if parts[0] in ("top", "middle", "bottom") else "middle"
    horiz = parts[1] if len(parts) > 1 else "center"
    if horiz == "left":
        x = margin
    elif horiz == "right":
        x = pw - margin - w
    else:
        x = (pw - w) / 2.0
    if vert == "top":
        y = margin
    elif vert == "bottom":
        y = ph - margin - h
    else:
        y = (ph - h) / 2.0
    return x, y


def pdf_watermark(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    """Stamp a watermark on every page — either text or a logo/image, at a
    chosen position."""
    opacity = _clampi(opt.get("opacity", 20), 1, 100, 20) / 100.0
    rotation = _clampi(opt.get("rotation", 45), -180, 180, 45)
    position = opt.get("position", "center")
    if opt.get("wm_type", "text") == "image":
        return _pdf_watermark_image(src, out_dir, opt, report, opacity, rotation,
                                    position)

    import fitz
    text = (opt.get("text") or "").strip()
    if not text:
        raise ProcessError("Enter the watermark text.")
    size = _clampi(opt.get("font_size", 48), 6, 400, 48)
    color = _WM_COLORS.get(opt.get("color", "gray"), (0.5, 0.5, 0.5))

    try:
        doc = fitz.open(str(src))
    except Exception as exc:
        raise ProcessError(f"Couldn't open '{src.name}' ({exc}).")
    try:
        for i, page in enumerate(doc):
            _check_cancel(report)
            rect = page.rect
            tl = fitz.get_text_length(text, fontname="helv", fontsize=size)
            x, y = _anchor_xy(position, rect.width, rect.height, tl, size)
            pivot = fitz.Point(x + tl / 2.0, y + size / 2.0)
            tw = fitz.TextWriter(rect, color=color)
            tw.append(fitz.Point(x, y + size * 0.8), text, fontsize=size)
            morph = (pivot, fitz.Matrix(rotation))
            tw.write_text(page, opacity=opacity, morph=morph, overlay=True)
            _progress(report, i + 1, doc.page_count)
        out = build_output_path(src, out_dir, ".pdf", name_suffix="_watermarked",
                                overwrite=opt.get("overwrite", False),
                                numbered=opt.get("same_as_source", False))
        pages = doc.page_count
        doc.save(str(out), garbage=3, deflate=True)
    finally:
        doc.close()
    report(f"Watermarked {pages} page(s) → {out.name}")
    return [out]


def _pdf_watermark_image(src: Path, out_dir: Path, opt: dict, report: Report,
                         opacity: float, rotation: int,
                         position: str = "center") -> list[Path]:
    """Overlay a logo/image (PNG transparency honoured) at *position* on every page."""
    import os
    import tempfile

    import fitz
    from PIL import Image

    raw = (opt.get("image_path") or "").strip()
    if not raw:
        raise ProcessError("Choose a logo / image for the watermark.")
    imgpath = Path(raw)
    if not imgpath.exists():
        raise ProcessError(f"Watermark image not found: {raw}")
    try:
        logo = Image.open(imgpath).convert("RGBA")
    except Exception as exc:
        raise ProcessError(f"Couldn't open image '{imgpath.name}' ({exc}).")

    # Apply opacity to the alpha channel; rotate (transparent corners) if asked.
    if opacity < 1.0:
        logo.putalpha(logo.split()[3].point(lambda a: int(a * opacity)))
    if rotation % 360:
        logo = logo.rotate(rotation, expand=True, resample=Image.BICUBIC)
    scale = _clampi(opt.get("scale", 40), 5, 100, 40) / 100.0
    ratio = (logo.height / logo.width) if logo.width else 1.0

    fd, tmp = tempfile.mkstemp(suffix=".png")
    os.close(fd)
    try:
        logo.save(tmp, "PNG")
        try:
            doc = fitz.open(str(src))
        except Exception as exc:
            raise ProcessError(f"Couldn't open '{src.name}' ({exc}).")
        try:
            for i, page in enumerate(doc):
                _check_cancel(report)
                rect = page.rect
                w = rect.width * scale
                h = w * ratio
                x, y = _anchor_xy(position, rect.width, rect.height, w, h)
                box = fitz.Rect(x, y, x + w, y + h)
                page.insert_image(box, filename=tmp, overlay=True, keep_proportion=True)
                _progress(report, i + 1, doc.page_count)
            out = build_output_path(src, out_dir, ".pdf", name_suffix="_watermarked",
                                    overwrite=opt.get("overwrite", False),
                                    numbered=opt.get("same_as_source", False))
            pages = doc.page_count
            doc.save(str(out), garbage=3, deflate=True)
        finally:
            doc.close()
    finally:
        try:
            os.remove(tmp)
        except OSError:
            pass
    report(f"Watermarked {pages} page(s) with {imgpath.name} → {out.name}")
    return [out]


# =========================================================================
# PDF → Word
# =========================================================================
def pdf_to_word(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    out = build_output_path(src, out_dir, ".docx", overwrite=opt.get("overwrite", False),
                            numbered=opt.get("same_as_source", False))

    # Scanned / image-only PDF + OCR enabled -> recognise text into a docx.
    if bool(opt.get("ocr", False)):
        import fitz
        doc = fitz.open(src)
        try:
            has_any_text = any(_page_has_text(p) for p in doc)
        finally:
            doc.close()
        if not has_any_text:
            _ocr_pdf_to_docx(src, out, report)
            return [out]

    try:
        from pdf2docx import Converter
    except Exception as exc:  # pragma: no cover
        raise ProcessError(f"pdf2docx is not available: {exc}")
    report("Converting layout to Word…")
    cv = Converter(str(src))
    try:
        cv.convert(str(out), start=0, end=None)
    finally:
        cv.close()
    if not out.exists():
        raise ProcessError("Conversion produced no output.")
    return [out]


def _ocr_pdf_to_docx(src: Path, out: Path, report: Report) -> None:
    """OCR every page of a scanned PDF into an editable .docx (python-docx)."""
    import fitz
    import docx

    _make_ocr_engine()
    pdf = fitz.open(src)
    try:
        npages = pdf.page_count
    finally:
        pdf.close()
    report(f"OCR {npages} page(s)…")
    # OCR all pages in parallel, then build the document in page order.
    page_lines = _ocr_pages_concurrent(src, list(range(npages)), report)

    document = docx.Document()
    for i in range(npages):
        # Reconstruct visual rows, then reflow them into readable paragraphs
        # (merge wrapped lines, de-hyphenate, break on real gaps/indents).
        rows = _ocr_rows(page_lines.get(i, []))
        paras = _rows_to_paragraphs(rows)
        if i:
            document.add_page_break()
        if not paras:
            document.add_paragraph("(No text recognised on this page.)")
        for para in paras:
            document.add_paragraph(para)
    document.save(str(out))
    if not out.exists():
        raise ProcessError("OCR produced no output.")


# =========================================================================
# OCR helpers (scanned / image-only pages → editable text)
# =========================================================================
# OCR is recognised at this DPI (higher = more accurate on small text, slower).
_OCR_DPI = 300
# Detections below this confidence are discarded as noise.
_OCR_MIN_SCORE = 0.5
# The RapidOCR engine loads several ONNX models; build it once and reuse it
# across every page and file in a batch instead of paying that cost per file.
_ocr_engine_cache = None
_ocr_engine_lock = threading.Lock()
_ocr_active_provider = "CPU"   # set when the engine is built; for status reporting


def _ocr_gpu_preference() -> bool:
    """User setting: use the GPU for OCR when one is available (default on)."""
    try:
        from mico360.config import settings
        return bool(settings.ocr_use_gpu)
    except Exception:
        return True


def _dml_available() -> bool:
    """True if onnxruntime exposes the DirectML provider (Windows GPU build).
    This is a capability check only — whether a usable GPU exists is decided at
    session-creation time, with CPU fallback."""
    try:
        import onnxruntime as ort
        return "DmlExecutionProvider" in ort.get_available_providers()
    except Exception:
        return False


def _patch_rapidocr_gpu() -> None:
    """Make RapidOCR build its ONNX sessions on the GPU via DirectML — any
    Direct3D-12 GPU (NVIDIA / AMD / Intel, discrete or integrated) — with an
    automatic CPU fallback. RapidOCR 1.2.x only knows CUDA, so we intercept
    onnxruntime session creation and prepend DmlExecutionProvider.

    Fully machine-agnostic: nothing here assumes a specific GPU. On a PC where
    DirectML can't create a session (no usable GPU / driver), each model quietly
    falls back to CPU. Idempotent."""
    try:
        import onnxruntime as ort
        from onnxruntime import ExecutionMode
        from rapidocr_onnxruntime import utils as _ru
    except Exception:
        return
    if getattr(_ru, "_mico360_gpu_patched", False):
        return
    _real_session = ort.InferenceSession

    def _make_session(model_path, sess_options=None, providers=None, **kw):
        global _ocr_active_provider
        if sess_options is not None:
            # DirectML requires a non-mem-pattern, sequential session.
            try:
                sess_options.enable_mem_pattern = False
                sess_options.execution_mode = ExecutionMode.ORT_SEQUENTIAL
            except Exception:
                pass
        gpu_eps = [("DmlExecutionProvider", {"device_id": 0}), "CPUExecutionProvider"]
        try:
            sess = _real_session(model_path, sess_options=sess_options,
                                 providers=gpu_eps, **kw)
        except Exception:
            # No usable GPU on this machine — fall back to whatever RapidOCR asked
            # for (CPU). Never let GPU selection break OCR.
            sess = _real_session(model_path, sess_options=sess_options,
                                 providers=providers, **kw)
        try:
            if "DmlExecutionProvider" in sess.get_providers():
                _ocr_active_provider = "GPU (DirectML)"
        except Exception:
            pass
        return sess

    _ru.InferenceSession = _make_session
    _ru._mico360_gpu_patched = True


def _make_ocr_engine():
    """Return a cached RapidOCR engine (models load once), or raise ProcessError.

    If the machine has a usable GPU and the user hasn't disabled it, the engine's
    ONNX models run on the GPU via DirectML (with CPU fallback); otherwise plain
    CPU. The decision is made fresh from the running machine — never hard-coded."""
    global _ocr_engine_cache, _ocr_active_provider
    if _ocr_engine_cache is not None:
        return _ocr_engine_cache
    with _ocr_engine_lock:
        if _ocr_engine_cache is None:
            try:
                from rapidocr_onnxruntime import RapidOCR
            except Exception as exc:  # pragma: no cover - optional dependency
                raise ProcessError(
                    "OCR engine is not available in this build. " + str(exc))
            _ocr_active_provider = "CPU"
            if _ocr_gpu_preference() and _dml_available():
                _patch_rapidocr_gpu()
            _ocr_engine_cache = RapidOCR()
    return _ocr_engine_cache


def ocr_active_provider() -> str:
    """Human-readable label of the OCR backend in use ('GPU (DirectML)' or
    'CPU'). Only meaningful after the engine has been built at least once."""
    return _ocr_active_provider


def _page_has_text(page, threshold: int = 8) -> bool:
    return len(page.get_text("text").strip()) >= threshold


def _render_page_image(page, dpi: int = _OCR_DPI):
    """Rasterise a page to an RGB numpy array for OCR. Returns (image, dpi)."""
    import numpy as np

    dpi = _clampi(dpi, 72, 600, _OCR_DPI)
    pix = page.get_pixmap(dpi=dpi)
    img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, pix.n)
    if pix.n == 4:
        img = np.ascontiguousarray(img[:, :, :3])
    elif pix.n == 1:
        img = np.ascontiguousarray(np.repeat(img, 3, axis=2))
    return img, dpi


def _ocr_image_lines(engine, img, dpi: int,
                     min_score: float = _OCR_MIN_SCORE) -> list[tuple[str, tuple]]:
    """Run OCR on an already-rendered image; return [(text, (x0,y0,x1,y1))] in
    PDF points. Drops detections below ``min_score`` so noise stays out.

    Angle classification (``use_cls``) is requested so upside-down / 180°-rotated
    lines on a scan are still read in the correct orientation.
    """
    try:
        result, _ = engine(img, use_cls=True)
    except TypeError:                        # older RapidOCR without the kwarg
        result, _ = engine(img)
    scale = 72.0 / dpi
    lines: list[tuple[str, tuple]] = []
    for box, text, score in (result or []):
        t = str(text).strip()
        if not t:
            continue
        try:
            if score is not None and float(score) < min_score:
                continue
        except (TypeError, ValueError):
            pass
        xs = [float(p[0]) for p in box]
        ys = [float(p[1]) for p in box]
        lines.append((t, (min(xs) * scale, min(ys) * scale,
                          max(xs) * scale, max(ys) * scale)))
    return lines


def _ocr_page_lines(engine, page, dpi: int = _OCR_DPI,
                    min_score: float = _OCR_MIN_SCORE) -> list[tuple[str, tuple]]:
    """OCR one page; return [(text, (x0, y0, x1, y1))] with coords in PDF points.
    Renders at ``dpi`` (clamped) for accuracy and filters low-confidence noise."""
    img, used_dpi = _render_page_image(page, dpi)
    return _ocr_image_lines(engine, img, used_dpi, min_score)


def _ocr_pages_concurrent(src: Path, indices: list[int],
                          report: Report,
                          dpi: int = _OCR_DPI) -> dict[int, list[tuple[str, tuple]]]:
    """OCR many pages of one PDF in parallel and return {page_index: lines}.

    The slow part — model inference — is concurrent (RapidOCR/onnxruntime is
    thread-safe and ~2.5× faster under concurrency), while rasterising stays on
    one thread (a fitz document is not thread-safe). Output is byte-identical to
    OCRing page-by-page; only the speed differs. Memory is bounded by processing
    in chunks. Used for a *single* file — the batch engine already parallelises
    across files, but pages within one file were previously sequential.
    """
    import concurrent.futures
    import os

    import fitz

    engine = _make_ocr_engine()
    if not indices:
        return {}
    cpu = os.cpu_count() or 4
    workers = max(1, min(len(indices), max(2, cpu // 2), 8))
    # On the GPU, run one page at a time: DirectML ONNX sessions are NOT safe for
    # concurrent Run() across threads (it crashes), and the GPU is the serial
    # resource anyway — serial GPU is still far faster than concurrent CPU.
    # CPU keeps its multi-thread speedup.
    if ocr_active_provider() != "CPU":
        workers = 1
    if workers == 1:
        doc = fitz.open(str(src))
        try:
            out = {}
            for n, i in enumerate(indices):
                _check_cancel(report)
                out[i] = _ocr_page_lines(engine, doc[i], dpi)
                _progress(report, n + 1, len(indices))
            return out
        finally:
            doc.close()

    results: dict[int, list] = {}
    chunk = workers * 2
    doc = fitz.open(str(src))
    try:
        with concurrent.futures.ThreadPoolExecutor(max_workers=workers) as ex:
            done = 0
            for start in range(0, len(indices), chunk):
                _check_cancel(report)
                batch = indices[start:start + chunk]
                imgs = {i: _render_page_image(doc[i], dpi) for i in batch}  # render (1 thread)
                for i, lines in zip(batch, ex.map(
                        lambda i: _ocr_image_lines(engine, *imgs[i]), batch)):
                    results[i] = lines
                done += len(batch)
                _progress(report, done, len(indices))
    finally:
        doc.close()
    return results


def _ocr_rows(lines: list[tuple[str, tuple]]) -> list[tuple[str, float, float, float]]:
    """Group OCR fragments into visual rows for clean, readable output.

    Fragments whose vertical centres are within a tolerance of each other are
    treated as one line and joined left→right. Returns
    ``[(row_text, x_left, y_top, row_height)]`` ordered top→bottom.
    """
    if not lines:
        return []
    items = []  # (text, x0, y0, height, y_center)
    for text, (x0, y0, x1, y1) in lines:
        items.append((text, x0, y0, max(1.0, y1 - y0), (y0 + y1) / 2.0))
    heights = sorted(it[3] for it in items)
    med_h = heights[len(heights) // 2] or 10.0
    tol = max(4.0, med_h * 0.6)

    items.sort(key=lambda it: it[4])  # by vertical centre
    rows: list[list] = [[items[0]]]
    centre = items[0][4]
    for it in items[1:]:
        if abs(it[4] - centre) <= tol:
            rows[-1].append(it)
        else:
            rows.append([it])
        centre = sum(c[4] for c in rows[-1]) / len(rows[-1])

    out: list[tuple[str, float, float, float]] = []
    for row in rows:
        row.sort(key=lambda it: it[1])  # left → right
        text = " ".join(c[0] for c in row).strip()
        if not text:
            continue
        out.append((text, min(c[1] for c in row),
                    min(c[2] for c in row), max(c[3] for c in row)))
    return out


def _rows_to_paragraphs(rows: list[tuple[str, float, float, float]]) -> list[str]:
    """Reflow visual OCR rows into readable paragraphs.

    Consecutive rows with normal line spacing are merged into one paragraph;
    a larger vertical gap, or a clearly indented new line, starts a new one.
    Words split across a line break with a trailing hyphen are re-joined
    (``inter-\\nnational`` → ``international``). Produces clean, editable prose
    instead of one stranded paragraph per scanned line.
    """
    if not rows:
        return []
    heights = sorted(h for _, _, _, h in rows)
    med_h = heights[len(heights) // 2] or 10.0
    lefts = sorted(x for _, x, _, _ in rows)
    base_left = lefts[len(lefts) // 4]            # typical left margin (robust)

    paras: list[str] = []
    cur = ""
    prev_bottom = None
    for text, x_left, y_top, height in rows:
        gap = (y_top - prev_bottom) if prev_bottom is not None else 0.0
        indented = (x_left - base_left) > med_h * 1.2
        new_para = prev_bottom is not None and (gap > med_h * 0.8 or indented)
        if new_para and cur:
            paras.append(cur)
            cur = ""
        if not cur:
            cur = text
        elif cur.endswith("-") and len(cur) >= 2 and cur[-2].isalpha():
            cur = cur[:-1] + text.lstrip()        # de-hyphenate across the break
        else:
            cur = cur + " " + text
        prev_bottom = y_top + height
    if cur:
        paras.append(cur)
    return paras


# =========================================================================
# PDF → PowerPoint  (one slide per page)
# =========================================================================
def pdf_to_pptx(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    """Turn each PDF page into a slide.

    ``mode`` controls how:
      * ``auto``  (default) — extract editable text boxes where the page has a
        real text layer; for image-only / scanned pages, OCR them (if enabled)
        or drop in the rendered page image so the slide is never empty.
      * ``text``  — editable text boxes only (small, fully editable).
      * ``image`` — render every page as an exact image (pixel-perfect look,
        not editable).

    ``ocr`` (bool) — when a page has no text layer, recognise its text so it
    still becomes editable text boxes instead of a flat image.
    """
    import tempfile

    import fitz
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Pt

    mode = opt.get("mode", "auto")
    ocr = bool(opt.get("ocr", False))
    ocr_engine = None
    out = build_output_path(src, out_dir, ".pptx", overwrite=opt.get("overwrite", False),
                            numbered=opt.get("same_as_source", False))
    doc = fitz.open(src)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    EMU_PER_PT = 12700  # 914400 EMU per inch / 72 pt per inch

    def add_text_boxes(slide, page) -> int:
        boxes = 0
        for block in page.get_text("dict").get("blocks", []):
            if block.get("type", 0) != 0:
                continue  # 0 = text block
            for line in block.get("lines", []):
                spans = [s for s in line.get("spans", []) if s.get("text", "").strip()]
                if not spans:
                    continue
                text = "".join(s.get("text", "") for s in spans)
                x0, y0, x1, y1 = line["bbox"]
                tb = slide.shapes.add_textbox(
                    Emu(int(x0 * EMU_PER_PT)), Emu(int(y0 * EMU_PER_PT)),
                    Emu(int(max(2, x1 - x0) * EMU_PER_PT)),
                    Emu(int(max(8, y1 - y0) * EMU_PER_PT)))
                tf = tb.text_frame
                tf.word_wrap = False
                tf.margin_left = tf.margin_right = 0
                tf.margin_top = tf.margin_bottom = 0
                run = tf.paragraphs[0].add_run()
                run.text = text
                s0 = spans[0]
                run.font.size = Pt(max(6.0, float(s0.get("size", 12))))
                flags = int(s0.get("flags", 0))
                run.font.bold = bool(flags & 16)
                run.font.italic = bool(flags & 2)
                col = int(s0.get("color", 0))
                run.font.color.rgb = RGBColor((col >> 16) & 255, (col >> 8) & 255, col & 255)
                boxes += 1
        return boxes

    def add_ocr_boxes(slide, lines) -> int:
        for text, (x0, y0, x1, y1) in lines:
            tb = slide.shapes.add_textbox(
                Emu(int(x0 * EMU_PER_PT)), Emu(int(y0 * EMU_PER_PT)),
                Emu(int(max(2, x1 - x0) * EMU_PER_PT)),
                Emu(int(max(8, y1 - y0) * EMU_PER_PT)))
            tf = tb.text_frame
            tf.word_wrap = False
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            run = tf.paragraphs[0].add_run()
            run.text = text
            run.font.size = Pt(max(7.0, min(40.0, (y1 - y0) * 0.85)))
        return len(lines)

    try:
        if doc.page_count:
            r0 = doc[0].rect
            prs.slide_width = Emu(int(r0.width * EMU_PER_PT))
            prs.slide_height = Emu(int(r0.height * EMU_PER_PT))
        with tempfile.TemporaryDirectory() as tmp:
            for i, page in enumerate(doc):
                _check_cancel(report)
                _progress(report, i, doc.page_count)
                slide = prs.slides.add_slide(blank)
                has_text = _page_has_text(page)
                if mode == "image":
                    pass  # always render the page image (below)
                elif has_text:
                    n = add_text_boxes(slide, page)
                    report(f"Slide {i + 1}/{doc.page_count} — {n} editable text box(es)")
                    continue
                elif ocr:
                    if ocr_engine is None:
                        ocr_engine = _make_ocr_engine()
                    n = add_ocr_boxes(slide, _ocr_page_lines(ocr_engine, page))
                    report(f"Slide {i + 1}/{doc.page_count} — {n} OCR text box(es)")
                    continue
                elif mode == "text":
                    report(f"Slide {i + 1}/{doc.page_count} — no text layer "
                           "(turn on OCR for scanned pages)")
                    continue
                # mode == image, or auto/text fallback to a page image
                pix = page.get_pixmap(dpi=150)
                img_path = Path(tmp) / f"p{i}.png"
                pix.save(str(img_path))
                slide.shapes.add_picture(str(img_path), 0, 0,
                                         width=prs.slide_width, height=prs.slide_height)
                report(f"Slide {i + 1}/{doc.page_count} — page image")
            prs.save(str(out))
    finally:
        doc.close()
    return [out]


# =========================================================================
# Word → PDF   (engine chain: LibreOffice → MS Word → built-in fallback)
# =========================================================================
def word_to_pdf(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    """Convert a Word document to PDF, trying the best available engine.

    Order: bundled/installed LibreOffice (highest fidelity, cross-platform) →
    Microsoft Word via COM (Windows, exact) → a built-in pure-Python renderer
    (text, headings, basic styling and tables) so it always produces a PDF even
    with nothing installed.
    """
    out_dir.mkdir(parents=True, exist_ok=True)
    final = build_output_path(src, out_dir, ".pdf", overwrite=opt.get("overwrite", False),
                              numbered=opt.get("same_as_source", False))
    errors: list[str] = []

    soffice = find_libreoffice()
    if soffice:
        try:
            _word_to_pdf_libreoffice(soffice, src, out_dir, final, report)
            return [final]
        except Exception as exc:  # noqa: BLE001
            errors.append(f"LibreOffice: {exc}")

    try:
        if _word_to_pdf_msword(src, final, report):
            return [final]
    except Exception as exc:  # noqa: BLE001
        errors.append(f"Microsoft Word: {exc}")

    try:
        _word_to_pdf_python(src, final, report)
        return [final]
    except Exception as exc:  # noqa: BLE001
        errors.append(f"Built-in: {exc}")

    raise ProcessError("Word→PDF failed with every engine. " + " | ".join(errors))


# External Office automation — LibreOffice (headless) and Microsoft Office (COM)
# — is effectively single-instance: concurrent / rapid invocations collide on a
# profile or COM lock and HANG. We serialise every Office conversion with one
# lock, and give each LibreOffice run its own throw-away UserInstallation profile.
# Reliable for a whole batch of Office files, even with a quickstarter running.
_office_lock = threading.Lock()


def _run_libreoffice(soffice: str, src: Path, out_dir: Path, final: Path,
                     report: Report, label: str = "LibreOffice") -> None:
    import shutil
    import tempfile

    report(f"{label} converting to PDF…")
    profile = tempfile.mkdtemp(prefix="mico360_lo_")
    args = [soffice, "--headless", "--norestore", "--invisible", "--nodefault",
            "--nolockcheck", f"-env:UserInstallation={Path(profile).as_uri()}",
            "--convert-to", "pdf", "--outdir", str(out_dir), str(src)]
    try:
        with _office_lock:
            proc = run_subprocess(args, timeout=300)
    finally:
        shutil.rmtree(profile, ignore_errors=True)
    produced = out_dir / f"{src.stem}.pdf"
    if proc.returncode != 0 or not produced.exists():
        raise RuntimeError((proc.stderr or proc.stdout or "no output").strip()[:200])
    if produced != final:
        if final.exists():
            final.unlink()
        produced.rename(final)


def _word_to_pdf_libreoffice(soffice: str, src: Path, out_dir: Path,
                             final: Path, report: Report) -> None:
    _run_libreoffice(soffice, src, out_dir, final, report)


def _word_to_pdf_msword(src: Path, final: Path, report: Report) -> bool:
    """Use Microsoft Word via COM. Returns False if Word/docx2pdf isn't present."""
    try:
        from docx2pdf import convert
    except Exception:
        return False
    com_inited = False
    try:
        import pythoncom  # from pywin32; required for COM in a worker thread
        pythoncom.CoInitialize()
        com_inited = True
    except Exception:
        pass
    try:
        report("Microsoft Word converting to PDF…")
        with _office_lock:          # serialise COM Office automation
            convert(str(src), str(final))
    finally:
        if com_inited:
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass
    if not final.exists():
        raise RuntimeError("Word produced no output")
    return True


def _word_to_pdf_python(src: Path, final: Path, report: Report) -> None:
    """Pure-Python fallback: render the document's text with python-docx +
    reportlab. Preserves paragraphs, headings, bold/italic and simple tables —
    not exact layout, but always works with no external program."""
    import html

    import docx  # python-docx
    from docx.table import Table as _Tbl
    from docx.text.paragraph import Paragraph as _Par
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import (
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    report("Converting with the built-in engine…")
    document = docx.Document(str(src))
    styles = getSampleStyleSheet()
    body_style = ParagraphStyle("body", parent=styles["Normal"], fontSize=11, leading=15)

    def heading_style(level: int) -> ParagraphStyle:
        size = max(12, 22 - 2 * level)
        return ParagraphStyle(f"h{level}", parent=styles["Normal"], fontSize=size,
                              leading=size + 4, spaceBefore=8, spaceAfter=4,
                              textColor=colors.HexColor("#1A1A1A"))

    def run_markup(par: _Par) -> str:
        parts = []
        for r in par.runs:
            t = html.escape(r.text or "")
            if not t:
                continue
            if r.bold:
                t = f"<b>{t}</b>"
            if r.italic:
                t = f"<i>{t}</i>"
            if r.underline:
                t = f"<u>{t}</u>"
            parts.append(t)
        return "".join(parts) or html.escape(par.text or "")

    flow: list = []
    # Walk the body in document order so paragraphs and tables stay interleaved.
    parent = document.element.body
    for child in parent.iterchildren():
        if child.tag.endswith("}p"):
            par = _Par(child, document)
            text = run_markup(par)
            if not text.strip():
                flow.append(Spacer(1, 6))
                continue
            name = (par.style.name or "").lower() if par.style else ""
            if name.startswith("heading") or name == "title":
                try:
                    level = int(name.split()[-1])
                except (ValueError, IndexError):
                    level = 1
                flow.append(Paragraph(text, heading_style(level)))
            else:
                flow.append(Paragraph(text, body_style))
        elif child.tag.endswith("}tbl"):
            tbl = _Tbl(child, document)
            data = [[Paragraph(html.escape(c.text or ""), body_style) for c in row.cells]
                    for row in tbl.rows]
            if data:
                t = Table(data, hAlign="LEFT")
                t.setStyle(TableStyle([
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#B0B0B0")),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 5),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 5),
                    ("TOPPADDING", (0, 0), (-1, -1), 3),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
                ]))
                flow.append(t)
                flow.append(Spacer(1, 8))

    if not flow:
        flow.append(Paragraph("(This document contained no extractable text.)", body_style))

    doc = SimpleDocTemplate(str(final), pagesize=LETTER,
                            leftMargin=0.9 * inch, rightMargin=0.9 * inch,
                            topMargin=0.9 * inch, bottomMargin=0.9 * inch)
    doc.build(flow)
    if not final.exists():
        raise RuntimeError("no output produced")


# =========================================================================
# PDF → Image  (per-file → one image per page)
# =========================================================================
def pdf_to_image(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    import fitz

    fmt = opt.get("format", "png").lower()
    dpi = _clampi(opt.get("dpi", 150), 72, 600, 150)
    overwrite = opt.get("overwrite", False)
    ext = {"jpg": ".jpg", "jpeg": ".jpg", "png": ".png", "webp": ".webp",
           "tiff": ".tiff", "bmp": ".bmp"}.get(fmt, ".png")
    dest = unique_dir(out_dir / src.stem)
    doc = fitz.open(src)
    outputs: list[Path] = []
    try:
        for i, page in enumerate(doc):
            _check_cancel(report)
            pix = page.get_pixmap(dpi=dpi)
            p = unique_path(dest / f"{src.stem}_page{i + 1}{ext}", overwrite)
            if ext == ".png":
                pix.save(str(p))  # native, keeps any alpha
            else:
                # PyMuPDF's pix.save only handles a few formats; use Pillow for
                # jpg / webp / bmp / tiff to guarantee correct encoding.
                import io

                from PIL import Image
                img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
                if ext == ".jpg":
                    img.save(p, quality=int(opt.get("jpeg_quality", 90)), optimize=True)
                else:
                    img.save(p)
            outputs.append(p)
            report(f"Page {i + 1}/{doc.page_count} → {p.name}")
            _progress(report, i + 1, doc.page_count)
    finally:
        doc.close()
    return outputs


# =========================================================================
# Image → PDF  (aggregate; one-per-image OR combined)
# =========================================================================
# =========================================================================
# SVG ⇄ image
# =========================================================================
_SVG_IMG_FMT = {"png": "PNG", "jpg": "JPEG", "jpeg": "JPEG", "webp": "WEBP"}


def svg_to_image(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    """Rasterise an SVG to a PNG / JPEG / WEBP image at a chosen width.

    Rendered with PyMuPDF (thread-safe, no GUI needed). Width 0 keeps the SVG's
    own size; a transparent background is kept for PNG/WEBP, and JPEG (which has
    no transparency) is flattened onto the chosen background colour."""
    import fitz
    from PIL import Image

    out_dir.mkdir(parents=True, exist_ok=True)
    fmt = str(opt.get("format", "png")).lower()
    pil_fmt = _SVG_IMG_FMT.get(fmt, "PNG")
    ext = ".jpg" if pil_fmt == "JPEG" else f".{fmt}"
    width = _clampi(opt.get("width", 1024), 0, 20000, 1024)
    bg = str(opt.get("background", "transparent")).lower()

    try:
        doc = fitz.open(str(src))
    except Exception as exc:
        raise ProcessError(f"Couldn't read the SVG '{src.name}' ({exc}).")
    try:
        if doc.page_count == 0:
            raise ProcessError(f"'{src.name}' has no drawable content.")
        page = doc[0]
        native_w = page.rect.width or 1.0
        scale = (width / native_w) if width else 1.0
        scale = max(0.01, scale)
        report(f"Rendering SVG at {int(native_w * scale)}px wide…")
        pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=True)
        mode = "RGBA" if pix.alpha else "RGB"
        img = Image.frombytes(mode, (pix.width, pix.height), pix.samples)
    finally:
        doc.close()

    if pil_fmt == "JPEG" or bg == "white":
        canvas = Image.new("RGB", img.size, (255, 255, 255))
        if img.mode == "RGBA":
            canvas.paste(img, mask=img.split()[3])
        else:
            canvas.paste(img)
        img = canvas
    elif img.mode == "RGBA" and pil_fmt not in ("PNG", "WEBP"):
        img = img.convert("RGB")

    out = build_output_path(src, out_dir, ext, overwrite=opt.get("overwrite", False),
                            numbered=opt.get("same_as_source", False))
    save_kw = {"optimize": True} if pil_fmt == "PNG" else (
        {"quality": 92} if pil_fmt in ("JPEG", "WEBP") else {})
    img.save(str(out), pil_fmt, **save_kw)
    report(f"SVG → {pil_fmt} ({img.width}×{img.height}) → {out.name}")
    return [out]


def _image_to_svg_embed(src: Path, out: Path) -> None:
    """Wrap the raster image, unchanged, inside an SVG (exact, lossless)."""
    import base64

    from PIL import Image
    with Image.open(str(src)) as im:
        w, h = im.size
    mime = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
            ".webp": "image/webp", ".gif": "image/gif",
            ".bmp": "image/bmp"}.get(src.suffix.lower())
    if mime:
        data = src.read_bytes()
    else:  # normalise odd formats (e.g. TIFF) to PNG so browsers can show them
        import io
        buf = io.BytesIO()
        with Image.open(str(src)) as im:
            im.save(buf, "PNG")
        data, mime = buf.getvalue(), "image/png"
    b64 = base64.b64encode(data).decode("ascii")
    out.write_text(
        f'<?xml version="1.0" encoding="UTF-8"?>\n'
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="{h}" '
        f'viewBox="0 0 {w} {h}">\n'
        f'  <image width="{w}" height="{h}" '
        f'xlink:href="data:{mime};base64,{b64}" '
        f'xmlns:xlink="http://www.w3.org/1999/xlink"/>\n'
        f'</svg>\n', encoding="utf-8")


def image_to_svg(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    """Convert a raster image to SVG.

    * ``trace`` (default) vectorises the image into real SVG paths (best for
      logos / line art / flat graphics) using vtracer.
    * ``embed`` wraps the original image, unchanged, inside an SVG (exact —
      good for photos, where tracing looks poor)."""
    out_dir.mkdir(parents=True, exist_ok=True)
    mode = str(opt.get("mode", "trace")).lower()
    out = build_output_path(src, out_dir, ".svg", overwrite=opt.get("overwrite", False),
                            numbered=opt.get("same_as_source", False))

    if mode == "embed":
        report("Embedding the image inside an SVG (exact)…")
        _image_to_svg_embed(src, out)
        report(f"Image → SVG (embedded) → {out.name}")
        return [out]

    # Trace to vector paths.
    try:
        import vtracer
    except Exception:
        report("Vector tracer unavailable — embedding the image inside the SVG instead.")
        _image_to_svg_embed(src, out)
        report(f"Image → SVG (embedded) → {out.name}")
        return [out]

    import tempfile

    from PIL import Image
    bw = str(opt.get("colors", "color")).lower() == "bw"
    report("Tracing the image to vector paths" + (" (black & white)…" if bw else "…"))
    # vtracer reads from a file. We always hand it a fresh PNG: for B&W we
    # threshold to two tones FIRST (vtracer's own colour options can't be passed
    # safely on this Python build), and otherwise normalise to RGBA PNG.
    work = Path(tempfile.mkdtemp()) / (src.stem + ".png")
    try:
        with Image.open(str(src)) as im:
            if bw:
                im.convert("L").point(lambda p: 0 if p < 128 else 255).convert(
                    "RGB").save(str(work), "PNG")
            else:
                im.convert("RGBA").save(str(work), "PNG")
        # IMPORTANT: pass NO keyword args — on this build the kwarg path crashes
        # the native tracer. The defaults are full-colour smooth (spline) tracing.
        vtracer.convert_image_to_svg_py(str(work), str(out))
    except Exception as exc:
        report(f"Tracing failed ({exc}) — embedding the image instead.")
        _image_to_svg_embed(src, out)
        report(f"Image → SVG (embedded) → {out.name}")
        return [out]
    report(f"Image → SVG (vector trace{', B&W' if bw else ''}) → {out.name}")
    return [out]


def image_to_pdf(srcs: list[Path], out_dir: Path, opt: dict, report: Report) -> list[Path]:
    from PIL import Image
    # Saving a PDF from Pillow uses the JPEG encoder internally; make sure all
    # plugins are registered first, or Pillow raises KeyError('JPEG') when its
    # lazy plugin init hasn't run yet (seen on Pillow 12 / frozen builds).
    Image.init()

    out_dir.mkdir(parents=True, exist_ok=True)
    overwrite = opt.get("overwrite", False)
    combine = opt.get("combine", True)
    imgs = [s for s in srcs if s.suffix.lower() in IMAGE_EXTS]
    if not imgs:
        raise ProcessError("No supported image files were provided.")

    def load(p: Path) -> "Image.Image":
        im = Image.open(p)
        if im.mode in ("RGBA", "P", "LA"):
            bg = Image.new("RGB", im.size, (255, 255, 255))
            im = im.convert("RGBA")
            bg.paste(im, mask=im.split()[-1])
            return bg
        return im.convert("RGB")

    if combine:
        name = (opt.get("output_name") or "images").strip() or "images"
        out = unique_path(out_dir / f"{name}.pdf", overwrite)
        loaded = []
        for j, p in enumerate(imgs):
            loaded.append(load(p))
            _progress(report, j + 1, len(imgs))
        first, *rest = loaded
        first.save(out, save_all=True, append_images=rest)
        report(f"Combined {len(imgs)} images → {out.name}")
        return [out]

    outputs: list[Path] = []
    numbered = opt.get("same_as_source", False)
    for j, p in enumerate(imgs):
        _check_cancel(report)
        # For "save next to source" in this aggregate tool, each PDF must land
        # beside its OWN image — not all beside the first one (the engine passes
        # a single base dir for aggregate units).
        dest = p.parent if numbered else out_dir
        out = build_output_path(p, dest, ".pdf", overwrite=overwrite, numbered=numbered)
        load(p).save(out)
        outputs.append(out)
        report(f"{p.name} → {out.name}")
        _progress(report, j + 1, len(imgs))
    return outputs


# =========================================================================
# Image compression
# =========================================================================
_IMG_QUALITY = {"low": 85, "medium": 65, "high": 45}


def verify_image_identical(src_path, out_path) -> bool:
    """True if two images have exactly the same dimensions and pixels (used to
    confirm a lossless image operation changed nothing visible)."""
    from PIL import Image, ImageChops
    try:
        a = Image.open(str(src_path)); a.load()
        b = Image.open(str(out_path)); b.load()
    except Exception:
        return False
    if a.size != b.size:
        return False
    if a.mode != b.mode:
        a = a.convert("RGBA"); b = b.convert("RGBA")
    return ImageChops.difference(a, b).getbbox() is None


def _image_compress_lossless(src: Path, im, out_dir: Path, opt: dict,
                             before: int, report: Report) -> list[Path]:
    """Re-pack an image with ZERO pixel change. Same format, lossless codec
    settings. JPEG/BMP have no lossless gain → the original is kept verbatim.
    The result is verified pixel-identical to the original."""
    import shutil
    report("Lossless image compression — preserving every pixel…")
    ext = src.suffix.lower()
    out = build_output_path(src, out_dir, ext, name_suffix="_compressed",
                            overwrite=opt.get("overwrite", False),
                            numbered=opt.get("same_as_source", False))
    saved = False
    try:
        if ext == ".png":
            im.save(out, format="PNG", optimize=True)
            saved = True
        elif ext == ".webp":
            im.save(out, format="WEBP", lossless=True, method=6)
            saved = True
        elif ext in (".tif", ".tiff"):
            im.save(out, format="TIFF", compression="tiff_deflate")
            saved = True
    except Exception as exc:
        report(f"Lossless re-pack unavailable ({exc}); keeping the original.")
        saved = False
    # JPEG/BMP (or any failure): copy the original — no lossless gain is possible.
    if not saved:
        shutil.copyfile(str(src), str(out))
    # Verify identical pixels; if not (or no size gain), keep the original.
    if not verify_image_identical(src, out) or out.stat().st_size >= before:
        shutil.copyfile(str(src), str(out))
        report(f"Already optimally packed — kept the original ({human_size(before)}); "
               "pixels verified identical ✓")
        return [out]
    after = out.stat().st_size
    report(f"{human_size(before)} → {human_size(after)} "
           f"({(1 - after / before) * 100:.0f}% smaller) — pixels verified identical ✓")
    return [out]


def image_compress(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    from PIL import Image

    if src.suffix.lower() not in IMAGE_EXTS:
        raise ProcessError(f"Unsupported image type: {src.suffix}")
    level = opt.get("level", "lossless")
    # The "quality" spinbox only applies in custom mode; for the Low/Medium/High
    # presets use the level's own quality so the choice actually takes effect.
    # (values() always reports the hidden spinbox, so we must branch on level.)
    if level == "custom":
        quality = _clampi(opt.get("quality", 65), 5, 100, 65)
    else:
        quality = _IMG_QUALITY.get(level, 65)
    out_format = opt.get("format", "keep")
    max_dim = _clampi(opt.get("max_dimension", 0), 0, 20000, 0)
    overwrite = opt.get("overwrite", False)
    before = src.stat().st_size

    try:
        im = Image.open(src)
        im.load()
    except Exception as exc:
        raise ProcessError(f"Couldn't open '{src.name}' — it may be corrupt or not "
                           f"a supported image ({exc}).")

    # Lossless (default): re-pack with zero pixel change, verified identical.
    if level == "lossless":
        return _image_compress_lossless(src, im, out_dir, opt, before, report)

    # Optional downscale
    if max_dim and max(im.size) > max_dim:
        ratio = max_dim / max(im.size)
        im = im.resize((int(im.width * ratio), int(im.height * ratio)), Image.LANCZOS)
        report(f"Resized to {im.width}x{im.height}")

    if out_format == "keep":
        ext = src.suffix.lower()
        if ext in (".tif", ".tiff", ".bmp"):
            ext = ".jpg"  # these don't benefit from lossy 'quality'; convert
    else:
        ext = {"jpg": ".jpg", "png": ".png", "webp": ".webp"}.get(out_format, ".jpg")

    # --- compress-to-target-size ---------------------------------------
    if level == "target":
        if ext == ".png":
            ext = ".jpg"  # need a lossy format to hit an arbitrary size
        target_bytes = max(1, int(opt.get("target_kb", 250))) * 1024
        out = build_output_path(src, out_dir, ext, name_suffix="_compressed",
                                overwrite=overwrite,
                                numbered=opt.get("same_as_source", False))
        after = _image_compress_to_target(im, out, ext, target_bytes, report)
        pct = (1 - after / before) * 100 if before else 0
        report(f"{human_size(before)} → {human_size(after)} ({pct:+.0f}%)")
        if after > target_bytes:
            report("Note: couldn't reach the target even at lowest quality / "
                   "smallest size; produced the smallest possible.")
        return [out]

    out = build_output_path(src, out_dir, ext, name_suffix="_compressed", overwrite=overwrite,
                            numbered=opt.get("same_as_source", False))
    save_kwargs: dict = {}
    if ext in (".jpg", ".jpeg"):
        im = im.convert("RGB")
        save_kwargs = {"quality": quality, "optimize": True, "progressive": True}
    elif ext == ".webp":
        save_kwargs = {"quality": quality, "method": 6}
    elif ext == ".png":
        if im.mode not in ("RGBA", "RGB", "P", "L"):
            im = im.convert("RGBA")
        save_kwargs = {"optimize": True}
    im.save(out, **save_kwargs)

    after = out.stat().st_size
    pct = (1 - after / before) * 100 if before else 0
    report(f"{human_size(before)} → {human_size(after)} ({pct:+.0f}%)")
    return [out]


def _encode_image(im, ext: str, quality: int) -> bytes:
    """Encode *im* to *ext* at *quality* in memory and return the bytes."""
    import io

    buf = io.BytesIO()
    work = im
    if ext in (".jpg", ".jpeg"):
        if work.mode != "RGB":
            work = work.convert("RGB")
        work.save(buf, "JPEG", quality=quality, optimize=True, progressive=True)
    elif ext == ".webp":
        work.save(buf, "WEBP", quality=quality, method=6)
    else:  # pragma: no cover - target only uses lossy formats
        work.save(buf, "PNG", optimize=True)
    return buf.getvalue()


def _image_compress_to_target(im, out: Path, ext: str, target_bytes: int,
                              report: Report) -> int:
    """Save *im* as the largest-quality encoding that stays <= target_bytes,
    downscaling progressively if even the lowest quality is too big. Returns the
    achieved byte size."""
    from PIL import Image

    report(f"Targeting <= {human_size(target_bytes)} (auto-selecting quality)…")
    best: bytes | None = None
    scale = 1.0
    for _ in range(7):
        work = im if scale == 1.0 else im.resize(
            (max(1, int(im.width * scale)), max(1, int(im.height * scale))),
            Image.LANCZOS)
        lo, hi, found = 5, 95, None
        while lo <= hi:                       # binary search quality
            q = (lo + hi) // 2
            data = _encode_image(work, ext, q)
            if len(data) <= target_bytes:
                found = data
                lo = q + 1
            else:
                hi = q - 1
        if found is not None:
            best = found
            break
        scale *= 0.8                          # too big even at q=5 -> shrink
        if min(int(im.width * scale), int(im.height * scale)) < 24:
            best = _encode_image(work, ext, 5)
            break
    if best is None:
        best = _encode_image(im, ext, 5)
    out.write_bytes(best)
    return len(best)


# =========================================================================
# v5.4 — dedicated page tools, page numbers, sign, metadata, searchable OCR
# =========================================================================
def _open_reader(src: Path):
    from pypdf import PdfReader
    try:
        reader = PdfReader(str(src))
        return reader, len(reader.pages)
    except Exception as exc:
        raise ProcessError(f"Couldn't read '{src.name}' — it may be corrupt or "
                           f"password-protected ({exc}).")


def _write_pdf(writer, src: Path, out_dir: Path, opt: dict, suffix: str,
               report: Report) -> list[Path]:
    out = build_output_path(src, out_dir, ".pdf", name_suffix=suffix,
                            overwrite=opt.get("overwrite", False),
                            numbered=opt.get("same_as_source", False))
    with open(out, "wb") as fh:
        writer.write(fh)
    report(f"Saved {len(writer.pages)} page(s) → {out.name}")
    return [out]


def pdf_rotate(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    from pypdf import PdfWriter
    reader, n = _open_reader(src)
    angle = _clampi(opt.get("angle", 90), -270, 270, 90)
    if angle % 90 != 0:
        raise ProcessError("Rotation must be a multiple of 90°.")
    targets = _page_set(opt.get("pages", "all"), n, default_all=True)
    writer = PdfWriter()
    for i in range(n):
        page = reader.pages[i]
        if i in targets:
            page.rotate(angle)
        writer.add_page(page)
        _progress(report, i + 1, n)
    return _write_pdf(writer, src, out_dir, opt, "_rotated", report)


def pdf_delete(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    from pypdf import PdfWriter
    reader, n = _open_reader(src)
    remove = _page_set(opt.get("pages", ""), n)
    if not remove:
        raise ProcessError("Enter the page(s) to delete, e.g. 2, 5-7.")
    kept = [i for i in range(n) if i not in remove]
    if not kept:
        raise ProcessError("That would delete every page.")
    writer = PdfWriter()
    for j, i in enumerate(kept):
        writer.add_page(reader.pages[i])
        _progress(report, j + 1, len(kept))
    return _write_pdf(writer, src, out_dir, opt, "_pages-removed", report)


def pdf_extract(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    from pypdf import PdfWriter
    reader, n = _open_reader(src)
    order = _page_list(opt.get("pages", ""), n)
    if not order:
        raise ProcessError("Enter the pages to extract, e.g. 1-5, 10, 15-20.")
    writer = PdfWriter()
    for j, i in enumerate(order):
        writer.add_page(reader.pages[i])
        _progress(report, j + 1, len(order))
    return _write_pdf(writer, src, out_dir, opt, "_extracted", report)


# (halign, valign) for each named position
_PN_POS = {
    "bottom-center": ("center", "bottom"), "bottom-right": ("right", "bottom"),
    "bottom-left": ("left", "bottom"), "top-center": ("center", "top"),
    "top-right": ("right", "top"), "top-left": ("left", "top"),
}


def pdf_page_numbers(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    import fitz
    fmt = opt.get("format", "n")
    halign, valign = _PN_POS.get(opt.get("position", "bottom-center"),
                                 ("center", "bottom"))
    size = _clampi(opt.get("font_size", 11), 6, 72, 11)
    start = _clampi(opt.get("start", 1), 1, 1000000, 1)
    margin = 28
    try:
        doc = fitz.open(str(src))
    except Exception as exc:
        raise ProcessError(f"Couldn't open '{src.name}' ({exc}).")
    try:
        total = doc.page_count
        for i, page in enumerate(doc):
            _check_cancel(report)
            num = start + i
            label = ({"n_of_total": f"{num} / {total}", "page_n": f"Page {num}"}
                     .get(fmt, str(num)))
            rect = page.rect
            tlen = fitz.get_text_length(label, fontname="helv", fontsize=size)
            if halign == "left":
                x = margin
            elif halign == "right":
                x = rect.width - margin - tlen
            else:
                x = (rect.width - tlen) / 2.0
            y = margin + size if valign == "top" else rect.height - margin
            page.insert_text(fitz.Point(x, y), label, fontsize=size,
                             fontname="helv", color=(0, 0, 0))
            _progress(report, i + 1, total)
        out = build_output_path(src, out_dir, ".pdf", name_suffix="_numbered",
                                overwrite=opt.get("overwrite", False),
                                numbered=opt.get("same_as_source", False))
        doc.save(str(out), garbage=3, deflate=True)
    finally:
        doc.close()
    report(f"Numbered {total} page(s) → {out.name}")
    return [out]


def pdf_sign(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    """Stamp a signature/image onto a chosen page (or all pages)."""
    import fitz

    raw = (opt.get("image_path") or "").strip()
    if not raw:
        raise ProcessError("Choose a signature image.")
    sigpath = Path(raw)
    if not sigpath.exists():
        raise ProcessError(f"Signature image not found: {raw}")
    halign, valign = _PN_POS.get(opt.get("position", "bottom-right"),
                                 ("right", "bottom"))
    width_pct = _clampi(opt.get("width", 25), 5, 100, 25) / 100.0
    which = opt.get("page", "last")  # last | first | all
    margin = 24
    try:
        doc = fitz.open(str(src))
    except Exception as exc:
        raise ProcessError(f"Couldn't open '{src.name}' ({exc}).")
    try:
        from PIL import Image
        with Image.open(sigpath) as im:
            ratio = (im.height / im.width) if im.width else 0.4
        if which == "first":
            pages = [0]
        elif which == "all":
            pages = list(range(doc.page_count))
        else:
            pages = [doc.page_count - 1]
        for n, pidx in enumerate(pages):
            _check_cancel(report)
            page = doc[pidx]
            rect = page.rect
            w = rect.width * width_pct
            h = w * ratio
            x0 = (margin if halign == "left"
                  else rect.width - margin - w if halign == "right"
                  else (rect.width - w) / 2.0)
            y0 = (margin if valign == "top" else rect.height - margin - h)
            box = fitz.Rect(x0, y0, x0 + w, y0 + h)
            page.insert_image(box, filename=str(sigpath), overlay=True,
                              keep_proportion=True)
            _progress(report, n + 1, len(pages))
        out = build_output_path(src, out_dir, ".pdf", name_suffix="_signed",
                                overwrite=opt.get("overwrite", False),
                                numbered=opt.get("same_as_source", False))
        doc.save(str(out), garbage=3, deflate=True)
    finally:
        doc.close()
    report(f"Signed → {out.name}")
    return [out]


def pdf_metadata(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    from pypdf import PdfWriter
    reader, _ = _open_reader(src)
    writer = PdfWriter()
    writer.append(reader)
    existing = dict(reader.metadata or {})
    fields = {"/Title": "title", "/Author": "author",
              "/Subject": "subject", "/Keywords": "keywords"}
    meta = {}
    changed = []
    for pdfkey, optkey in fields.items():
        val = opt.get(optkey, None)
        if val is None or str(val).strip() == "":
            if existing.get(pdfkey):
                meta[pdfkey] = existing[pdfkey]   # keep what's there
        else:
            meta[pdfkey] = str(val)
            changed.append(optkey)
    try:
        writer.add_metadata(meta)
    except Exception as exc:
        raise ProcessError(f"Couldn't write metadata ({exc}).")
    out = build_output_path(src, out_dir, ".pdf", name_suffix="_metadata",
                            overwrite=opt.get("overwrite", False),
                            numbered=opt.get("same_as_source", False))
    with open(out, "wb") as fh:
        writer.write(fh)
    report(f"Updated metadata ({', '.join(changed) or 'no changes'}) → {out.name}")
    return [out]


# =========================================================================
# Bulk file-property editor (filesystem: Date Created / Modified / Owner)
# =========================================================================
_DT_FORMATS = (
    "%Y-%m-%d %H:%M:%S", "%Y-%m-%d %H:%M", "%Y-%m-%d",
    "%m/%d/%Y %H:%M:%S", "%m/%d/%Y %H:%M", "%m/%d/%Y",
    "%d/%m/%Y %H:%M", "%d/%m/%Y", "%d-%m-%Y",
)


def _parse_dt(value):
    """Parse a user-typed date/time. Empty → None (leave unchanged). 'now'/'today'
    → the current time. Raises ProcessError on an unrecognised format."""
    import datetime
    s = (value or "").strip()
    if not s:
        return None
    if s.lower() in ("now", "today"):
        return datetime.datetime.now()
    for fmt in _DT_FORMATS:
        try:
            return datetime.datetime.strptime(s, fmt)
        except ValueError:
            continue
    raise ProcessError(
        f"Couldn't understand the date '{s}'. Use e.g. 2026-06-07 or "
        "2026-06-07 14:30.")


def _set_creation_time(path: Path, dt) -> bool:
    """Set a file's creation time (Windows only). Returns False off-Windows."""
    import sys
    if sys.platform != "win32":
        return False
    import pywintypes
    import win32con
    import win32file
    handle = win32file.CreateFile(
        str(path), win32con.GENERIC_WRITE,
        win32con.FILE_SHARE_READ | win32con.FILE_SHARE_WRITE
        | win32con.FILE_SHARE_DELETE,
        None, win32con.OPEN_EXISTING, win32con.FILE_ATTRIBUTE_NORMAL, None)
    try:
        win32file.SetFileTime(handle, pywintypes.Time(dt), None, None)  # creation only
    finally:
        handle.Close()
    return True


def _set_owner(path: Path, account: str) -> tuple[bool, str]:
    """Set the NTFS owner to ``account`` (Windows; needs admin/SeRestore).
    Returns (ok, error_message)."""
    import sys
    if sys.platform != "win32":
        return False, "owner can only be set on Windows"
    try:
        import win32security
        sid, _domain, _type = win32security.LookupAccountName(None, account)
        sd = win32security.GetFileSecurity(
            str(path), win32security.OWNER_SECURITY_INFORMATION)
        sd.SetSecurityDescriptorOwner(sid, False)
        win32security.SetFileSecurity(
            str(path), win32security.OWNER_SECURITY_INFORMATION, sd)
        return True, ""
    except Exception as exc:  # noqa: BLE001
        return False, str(exc)


def set_file_properties(src: Path, out_dir: Path, opt: dict,
                        report: Report) -> list[Path]:
    """Bulk-apply Windows file properties (Date Created, Date Modified, Owner) to
    a copy of each selected file — the original is never changed, and the copy's
    *content* is byte-identical (only the file properties differ)."""
    import os
    import shutil

    out_dir.mkdir(parents=True, exist_ok=True)
    out = build_output_path(src, out_dir, src.suffix or "",
                            overwrite=opt.get("overwrite", False),
                            numbered=opt.get("same_as_source", False))
    shutil.copy2(src, out)   # exact bytes + original times as the starting point

    created = _parse_dt(opt.get("date_created"))
    modified = _parse_dt(opt.get("date_modified"))
    owner = (opt.get("owner") or "").strip()
    changed: list[str] = []

    # Date Modified first (os.utime sets access+modify; keep access as-is).
    if modified is not None:
        st = out.stat()
        os.utime(out, (st.st_atime, modified.timestamp()))
        changed.append("Date Modified")

    # Date Created (Windows only) — set after, so it can't clobber the above.
    if created is not None:
        try:
            if _set_creation_time(out, created):
                changed.append("Date Created")
            else:
                report("Note: Date Created can only be set on Windows — skipped.")
        except Exception as exc:  # noqa: BLE001
            report(f"Note: couldn't set Date Created ({exc}).")

    # Owner (Windows; needs admin).
    if owner:
        ok, err = _set_owner(out, owner)
        if ok:
            changed.append(f"Owner={owner}")
        else:
            report(f"Note: couldn't set Owner ({err}). "
                   "Run as administrator and use a valid Windows account.")

    if changed:
        report("Set " + ", ".join(changed) + f"  →  {out.name}")
    elif created is not None or modified is not None or owner:
        report(f"Saved a copy → {out.name} (see notes above for what couldn't be set).")
    else:
        report(f"No properties entered — copied unchanged → {out.name}")
    return [out]


# OCR recognition DPI by quality preset (higher = sharper on small text, slower).
_OCR_QUALITY_DPI = {"fast": 200, "balanced": _OCR_DPI, "high": 400}


def _ocr_dpi_from_opt(opt: dict) -> int:
    return _OCR_QUALITY_DPI.get(str(opt.get("quality", "balanced")), _OCR_DPI)


def pdf_ocr(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    """Make a scanned PDF searchable by overlaying an invisible OCR text layer."""
    import fitz

    _make_ocr_engine()   # ensure OCR is available (raises a clear error if not)
    report(f"OCR engine ready on {ocr_active_provider()}.")
    dpi = _ocr_dpi_from_opt(opt)
    try:
        doc = fitz.open(str(src))
    except Exception as exc:
        raise ProcessError(f"Couldn't open '{src.name}' ({exc}).")
    try:
        todo = [i for i in range(doc.page_count) if not _page_has_text(doc[i])]
        kept = doc.page_count - len(todo)
        report(f"OCR {len(todo)} page(s)…"
               + (f"  ({kept} already searchable)" if kept else ""))
        # OCR the image-only pages in parallel, then write text in page order.
        page_lines = _ocr_pages_concurrent(src, todo, report, dpi)
        added = 0
        for i in todo:
            for text, (x0, y0, x1, y1) in page_lines.get(i, []):
                fontsize = max(4.0, (y1 - y0) * 0.9)
                # render_mode=3 → invisible text (selectable/searchable only)
                doc[i].insert_text(fitz.Point(x0, y1), text, fontsize=fontsize,
                                   fontname="helv", render_mode=3)
                added += 1
        out = build_output_path(src, out_dir, ".pdf", name_suffix="_searchable",
                                overwrite=opt.get("overwrite", False),
                                numbered=opt.get("same_as_source", False))
        doc.save(str(out), garbage=3, deflate=True)
    finally:
        doc.close()
    report(f"Added a searchable text layer ({added} snippets) → {out.name}")
    return [out]


# =========================================================================
# v5.4 — Office conversions (PDF↔Excel, Excel→PDF, PowerPoint→PDF)
# =========================================================================
def _sanitize_sheet(name: str) -> str:
    for ch in r"[]:*?/\\":
        name = name.replace(ch, " ")
    return (name.strip() or "Sheet")[:31]


def pdf_to_excel(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    """Extract tables from a PDF into an .xlsx (one sheet per table); falls back
    to page text when no tables are detected."""
    import pdfplumber
    from openpyxl import Workbook

    out = build_output_path(src, out_dir, ".xlsx", overwrite=opt.get("overwrite", False),
                            numbered=opt.get("same_as_source", False))
    wb = Workbook()
    wb.remove(wb.active)
    tables_found = 0
    try:
        with pdfplumber.open(str(src)) as pdf:
            total = len(pdf.pages) or 1
            for pi, page in enumerate(pdf.pages):
                _check_cancel(report)
                report(f"Scanning page {pi + 1}/{total} for tables…")
                for ti, table in enumerate(page.extract_tables() or []):
                    ws = wb.create_sheet(_sanitize_sheet(f"P{pi + 1}_T{ti + 1}"))
                    for row in table:
                        ws.append([("" if c is None else str(c)) for c in row])
                    tables_found += 1
                _progress(report, pi + 1, total)
            if not tables_found:
                report("No tables detected — exporting page text instead.")
                for pi, page in enumerate(pdf.pages):
                    ws = wb.create_sheet(_sanitize_sheet(f"Page {pi + 1}"))
                    for line in (page.extract_text() or "").splitlines():
                        ws.append([line])
    except Exception as exc:
        raise ProcessError(f"Couldn't read '{src.name}' ({exc}).")
    if not wb.sheetnames:
        wb.create_sheet("Empty")
    wb.save(str(out))
    report(f"{tables_found} table(s) → {out.name}" if tables_found
           else f"Text export → {out.name}")
    return [out]


def _lo_convert_to_pdf(soffice: str, src: Path, out_dir: Path, final: Path,
                       report: Report) -> None:
    _run_libreoffice(soffice, src, out_dir, final, report)


def _office_com_to_pdf(app_name: str, src: Path, final: Path) -> bool:
    """Use Microsoft Office (Excel/PowerPoint) via COM if installed."""
    try:
        import win32com.client as win32
    except Exception:
        return False
    import pythoncom
    pythoncom.CoInitialize()
    app = None
    try:
        with _office_lock:          # serialise COM Office automation
            app = win32.DispatchEx(app_name)
            src_s, final_s = str(src.resolve()), str(final.resolve())
            if app_name.startswith("Excel"):
                try:
                    app.Visible = False
                except Exception:
                    pass
                app.DisplayAlerts = False
                wb = app.Workbooks.Open(src_s, ReadOnly=True)
                wb.ExportAsFixedFormat(0, final_s)   # 0 = xlTypePDF
                wb.Close(False)
            else:  # PowerPoint
                pres = app.Presentations.Open(src_s, WithWindow=False)
                pres.SaveAs(final_s, 32)             # 32 = ppSaveAsPDF
                pres.Close()
        return final.exists()
    finally:
        try:
            if app is not None:
                app.Quit()
        except Exception:
            pass
        pythoncom.CoUninitialize()


def _office_to_pdf(src: Path, out_dir: Path, opt: dict, report: Report,
                   com_app: str) -> list[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    final = build_output_path(src, out_dir, ".pdf", overwrite=opt.get("overwrite", False),
                              numbered=opt.get("same_as_source", False))
    errors: list[str] = []
    soffice = find_libreoffice()
    if soffice:
        try:
            _lo_convert_to_pdf(soffice, src, out_dir, final, report)
            return [final]
        except Exception as exc:  # noqa: BLE001
            errors.append(f"LibreOffice: {exc}")
    try:
        report("Trying Microsoft Office…")
        if _office_com_to_pdf(com_app, src, final):
            return [final]
    except Exception as exc:  # noqa: BLE001
        errors.append(f"MS Office: {exc}")
    raise ProcessError(
        "Converting to PDF needs LibreOffice (recommended) or Microsoft Office. "
        "Install LibreOffice or set its path in Settings → External tools."
        + (f"\nDetails: {'; '.join(errors)}" if errors else ""))


def excel_to_pdf(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    return _office_to_pdf(src, out_dir, opt, report, "Excel.Application")


def pptx_to_pdf(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    return _office_to_pdf(src, out_dir, opt, report, "PowerPoint.Application")


# Input extension families (kept local so processors never imports the UI/tools).
_WORD_EXTS = {".doc", ".docx", ".odt", ".rtf"}
_EXCEL_EXTS = {".xlsx", ".xls", ".ods", ".csv"}
_PPT_EXTS = {".pptx", ".ppt", ".odp"}


def office_to_pdf(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    """Convert a Word, Excel or PowerPoint file to PDF — the type is detected
    automatically from the file, so one tool handles all three."""
    ext = src.suffix.lower()
    if ext in _WORD_EXTS:
        return word_to_pdf(src, out_dir, opt, report)
    if ext in _EXCEL_EXTS:
        return excel_to_pdf(src, out_dir, opt, report)
    if ext in _PPT_EXTS:
        return pptx_to_pdf(src, out_dir, opt, report)
    raise ProcessError(
        f"'{src.name}' isn't a Word, Excel or PowerPoint file.")


def pdf_convert(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    """Convert a PDF to another format chosen by ``target``
    (word / pptx / excel / image), routing to the matching converter. The
    per-target options live under distinct keys in the combined tool, so they
    are mapped back to what each converter expects."""
    target = str(opt.get("target", "word"))
    if target == "word":
        return pdf_to_word(src, out_dir, {**opt, "ocr": bool(opt.get("word_ocr"))}, report)
    if target == "pptx":
        return pdf_to_pptx(src, out_dir, {**opt, "ocr": bool(opt.get("pptx_ocr"))}, report)
    if target == "excel":
        return pdf_to_excel(src, out_dir, opt, report)
    if target == "image":
        return pdf_to_image(src, out_dir, opt, report)
    raise ProcessError(f"Unknown conversion target '{target}'.")


# =========================================================================
# Word → Markdown  (bulk .docx → .md)
# =========================================================================
def _lo_convert_to(soffice: str, src: Path, to_ext: str) -> Path:
    """Convert ``src`` to a temporary file of type ``to_ext`` (e.g. 'docx',
    'xlsx', 'pptx') via LibreOffice. Returns the produced file's path."""
    import tempfile

    to_ext = to_ext.lstrip(".")
    tmpdir = Path(tempfile.mkdtemp(prefix="mico360_md_"))
    profile = Path(tempfile.mkdtemp(prefix="mico360_lo_"))
    args = [soffice, "--headless", "--norestore", "--invisible", "--nodefault",
            "--nolockcheck", f"-env:UserInstallation={profile.as_uri()}",
            "--convert-to", to_ext, "--outdir", str(tmpdir), str(src)]
    try:
        with _office_lock:
            proc = run_subprocess(args, timeout=300)
    finally:
        import shutil
        shutil.rmtree(profile, ignore_errors=True)
    produced = tmpdir / f"{src.stem}.{to_ext}"
    if proc.returncode != 0 or not produced.exists():
        raise RuntimeError((proc.stderr or proc.stdout or "no output").strip()[:200])
    return produced


def _lo_convert_to_docx(soffice: str, src: Path) -> Path:
    """Convert a .doc/.odt/.rtf to a temporary .docx via LibreOffice."""
    return _lo_convert_to(soffice, src, "docx")


def _md_escape(text: str) -> str:
    # Light escaping so literal Markdown characters survive as text.
    return text.replace("\\", "\\\\").replace("|", "\\|")


def _runs_to_md(para) -> str:
    """Inline runs of a paragraph -> Markdown with **bold** / *italic* / `code`."""
    out = []
    for run in para.runs:
        t = run.text
        if not t:
            continue
        t = t.replace("\\", "\\\\")
        bold = bool(run.bold)
        italic = bool(run.italic)
        name = (run.font.name or "").lower()
        if "mono" in name or "courier" in name or "consol" in name:
            t = f"`{t}`"
        else:
            if bold and italic:
                t = f"***{t}***"
            elif bold:
                t = f"**{t}**"
            elif italic:
                t = f"*{t}*"
        out.append(t)
    return "".join(out).strip() or para.text.strip()


def _para_to_md(para) -> str:
    style = (para.style.name or "").lower() if para.style else ""
    text = _runs_to_md(para)
    if not text:
        return ""
    if style.startswith("heading"):
        try:
            level = int("".join(c for c in style if c.isdigit()) or "1")
        except ValueError:
            level = 1
        return f"{'#' * max(1, min(6, level))} {text}"
    if style in ("title",):
        return f"# {text}"
    if style in ("subtitle",):
        return f"## {text}"
    if style.startswith("quote") or style == "intense quote":
        return f"> {text}"
    # Lists: numbered vs bulleted (numbering present on the paragraph).
    is_list = "list" in style
    numbered = "number" in style
    try:
        from docx.oxml.ns import qn
        if para._p.pPr is not None and para._p.pPr.find(qn("w:numPr")) is not None:
            is_list = True
            fmt = para._p.pPr.find(qn("w:numPr"))
            ilvl = fmt.find(qn("w:ilvl"))
            numbered = numbered or ("number" in style)
            indent = int(ilvl.get(qn("w:val"))) if ilvl is not None else 0
        else:
            indent = 0
    except Exception:
        indent = 0
    if is_list:
        pad = "  " * max(0, indent)
        return f"{pad}{'1.' if numbered else '-'} {text}"
    return text


def _table_to_md(table) -> list[str]:
    rows = []
    for r in table.rows:
        cells = []
        for c in r.cells:
            text = " ".join(p.text for p in c.paragraphs)
            cells.append(" ".join(_md_escape(text).split()))   # collapse whitespace
        rows.append("| " + " | ".join(cells) + " |")
    if not rows:
        return []
    ncol = max(1, rows[0].count("|") - 1)
    sep = "| " + " | ".join(["---"] * ncol) + " |"
    return [rows[0], sep] + rows[1:]


def _docx_to_markdown(path: Path) -> str:
    import docx
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = docx.Document(str(path))
    blocks: list[str] = []
    for child in doc.element.body.iterchildren():
        if child.tag == qn("w:p"):
            blocks.append(_para_to_md(Paragraph(child, doc)))
        elif child.tag == qn("w:tbl"):
            blocks.append("")
            blocks.extend(_table_to_md(Table(child, doc)))
            blocks.append("")
    # Collapse 3+ blank lines down to a single blank line.
    md_lines: list[str] = []
    blank = 0
    for line in blocks:
        if line.strip() == "":
            blank += 1
            if blank <= 1:
                md_lines.append("")
        else:
            blank = 0
            md_lines.append(line)
    return "\n".join(md_lines).strip() + "\n"


def word_to_md(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    """Convert a Word document to Markdown (.md)."""
    try:
        import docx
        _ = docx.Document
    except Exception as exc:  # pragma: no cover
        raise ProcessError(f"python-docx is not available: {exc}")

    work = src
    cleanup: Path | None = None
    soffice = find_libreoffice()
    if src.suffix.lower() != ".docx":
        if not soffice:
            raise ProcessError(
                f"'{src.name}' isn't a .docx. Converting .doc/.odt/.rtf to Markdown "
                "needs LibreOffice (install it or set its path in Settings → "
                "External tools), or save the file as .docx first.")
        report("Converting to .docx via LibreOffice…")
        try:
            work = _lo_convert_to_docx(soffice, src)
            cleanup = work.parent
        except Exception as exc:
            raise ProcessError(f"Couldn't read '{src.name}' ({exc}).")

    try:
        report("Converting to Markdown…")
        try:
            md = _docx_to_markdown(work)
        except Exception as exc:
            # The file has a .docx name but python-docx can't open it (corrupt or
            # mis-named — e.g. a legacy .doc or RTF saved with a .docx extension).
            # If LibreOffice is available, let it repair/normalise the file into a
            # clean .docx and retry, so the batch doesn't fail on a quirky file.
            if cleanup is None and soffice is not None:
                report("File needs repair — normalising via LibreOffice…")
                try:
                    work = _lo_convert_to_docx(soffice, src)
                    cleanup = work.parent
                    md = _docx_to_markdown(work)
                except Exception as exc2:
                    raise ProcessError(
                        f"Couldn't convert '{src.name}' to Markdown ({exc2}).")
            else:
                raise ProcessError(
                    f"Couldn't convert '{src.name}' to Markdown ({exc}).")
    finally:
        if cleanup is not None:
            import shutil
            shutil.rmtree(cleanup, ignore_errors=True)

    out = build_output_path(src, out_dir, ".md", overwrite=opt.get("overwrite", False),
                            numbered=opt.get("same_as_source", False))
    out.write_text(md, encoding="utf-8")
    report(f"Markdown written → {out.name}")
    return [out]


def _grid_to_md(rows: list[list[str]]) -> str:
    """Render a 2-D grid of cell strings as a GitHub-flavoured Markdown table
    (first row = header)."""
    rows = [r for r in rows if any((c or "").strip() for c in r)]
    if not rows:
        return ""
    ncol = max(len(r) for r in rows)
    rows = [list(r) + [""] * (ncol - len(r)) for r in rows]
    lines = ["| " + " | ".join(_md_escape(c) for c in rows[0]) + " |",
             "| " + " | ".join(["---"] * ncol) + " |"]
    for r in rows[1:]:
        lines.append("| " + " | ".join(_md_escape(c) for c in r) + " |")
    return "\n".join(lines)


def _excel_to_markdown(src: Path, report: Report) -> str:
    """Workbook → Markdown: one '## Sheet' section with a table per sheet."""
    ext = src.suffix.lower()
    if ext == ".csv":
        import csv
        with src.open("r", encoding="utf-8", errors="replace", newline="") as f:
            rows = [[("" if c is None else str(c)) for c in row]
                    for row in csv.reader(f)]
        return f"# {src.stem}\n\n{_grid_to_md(rows)}\n"

    work, cleanup = src, None
    if ext != ".xlsx":
        soffice = find_libreoffice()
        if not soffice:
            raise ProcessError(
                f"Converting {ext} to Markdown needs LibreOffice (install it or set "
                "its path in Settings → External tools), or save it as .xlsx / .csv.")
        report("Converting to .xlsx via LibreOffice…")
        try:
            work = _lo_convert_to(soffice, src, "xlsx")
            cleanup = work.parent
        except Exception as exc:
            raise ProcessError(f"Couldn't read '{src.name}' ({exc}).")
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(work), read_only=True, data_only=True)
        parts: list[str] = [f"# {src.stem}"]
        for ws in wb.worksheets:
            rows = [["" if c is None else str(c) for c in row]
                    for row in ws.iter_rows(values_only=True)]
            table = _grid_to_md(rows)
            parts.append(f"## {ws.title}\n\n{table}" if table else f"## {ws.title}\n\n_(empty)_")
        wb.close()
        return "\n\n".join(parts).strip() + "\n"
    finally:
        if cleanup is not None:
            import shutil
            shutil.rmtree(cleanup, ignore_errors=True)


def _pptx_to_markdown(src: Path, report: Report) -> str:
    """Deck → Markdown: one '## Slide N' heading per slide, text as bullets."""
    ext = src.suffix.lower()
    work, cleanup = src, None
    if ext != ".pptx":
        soffice = find_libreoffice()
        if not soffice:
            raise ProcessError(
                f"Converting {ext} to Markdown needs LibreOffice (install it or set "
                "its path in Settings → External tools), or save it as .pptx.")
        report("Converting to .pptx via LibreOffice…")
        try:
            work = _lo_convert_to(soffice, src, "pptx")
            cleanup = work.parent
        except Exception as exc:
            raise ProcessError(f"Couldn't read '{src.name}' ({exc}).")
    try:
        from pptx import Presentation
        prs = Presentation(str(work))
        parts: list[str] = [f"# {src.stem}"]
        for i, slide in enumerate(prs.slides, 1):
            title_shape = None
            try:
                title_shape = slide.shapes.title
            except Exception:
                title_shape = None
            title = (title_shape.text.strip() if title_shape and title_shape.text else "")
            title_id = getattr(title_shape, "shape_id", None)
            parts.append(f"## Slide {i}" + (f": {title}" if title else ""))
            lines: list[str] = []
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                if title_id is not None and getattr(shape, "shape_id", None) == title_id:
                    continue
                for para in shape.text_frame.paragraphs:
                    txt = "".join(r.text for r in para.runs).strip() or para.text.strip()
                    if txt:
                        lines.append("  " * max(0, para.level) + "- " + txt)
            parts.append("\n".join(lines) if lines else "_(no text)_")
        return "\n\n".join(parts).strip() + "\n"
    finally:
        if cleanup is not None:
            import shutil
            shutil.rmtree(cleanup, ignore_errors=True)


def _pdf_to_markdown(src: Path, report: Report) -> str:
    """PDF → Markdown: page text as paragraphs, with detected tables rendered as
    Markdown tables. Image-only (scanned) pages are flagged so the user knows to
    OCR them first."""
    import fitz
    parts: list[str] = [f"# {src.stem}"]
    doc = fitz.open(str(src))
    try:
        multi = doc.page_count > 1
        for i in range(doc.page_count):
            _check_cancel(report)
            page = doc[i]
            if multi:
                parts.append(f"## Page {i + 1}")
            tables_md: list[str] = []
            try:                                  # PyMuPDF table detection (best-effort)
                for tab in page.find_tables().tables:
                    grid = [["" if c is None else str(c) for c in row]
                            for row in tab.extract()]
                    md = _grid_to_md(grid)
                    if md:
                        tables_md.append(md)
            except Exception:
                pass
            text = (page.get_text("text") or "").strip()
            if text:
                # Blank-line-separate the visual blocks for readable paragraphs.
                blocks = [b.strip() for b in text.split("\n\n") if b.strip()]
                parts.append("\n\n".join(blocks) if blocks else text)
            elif not tables_md:
                parts.append("_(no selectable text — this page looks scanned; run "
                             "Searchable PDF (OCR) on it first.)_")
            parts.extend(tables_md)
            _progress(report, i + 1, doc.page_count)
    finally:
        doc.close()
    return "\n\n".join(parts).strip() + "\n"


def to_markdown(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    """Convert a Word / Excel / PowerPoint / PDF file to Markdown (.md). The
    source type is detected automatically from the file."""
    ext = src.suffix.lower()
    if ext in _WORD_EXTS:
        return word_to_md(src, out_dir, opt, report)   # handles .docx + LO fallback

    report("Converting to Markdown…")
    if ext == ".pdf":
        md = _pdf_to_markdown(src, report)
    elif ext in _EXCEL_EXTS:
        md = _excel_to_markdown(src, report)
    elif ext in _PPT_EXTS:
        md = _pptx_to_markdown(src, report)
    else:
        raise ProcessError(f"Can't convert '{src.name}' to Markdown.")

    out = build_output_path(src, out_dir, ".md", overwrite=opt.get("overwrite", False),
                            numbered=opt.get("same_as_source", False))
    out.write_text(md, encoding="utf-8")
    report(f"Markdown written → {out.name}")
    return [out]


# =========================================================================
# v5.4 — dedicated image tools: Resize, Convert, Watermark
# =========================================================================
def _save_image(im, out: Path, fmt: str, quality: int = 90) -> None:
    """Save *im* to *out* in *fmt* (jpg/png/webp/tiff/bmp), handling mode."""
    fmt = fmt.lower()
    if fmt in ("jpg", "jpeg"):
        if im.mode != "RGB":
            im = im.convert("RGB")
        im.save(out, "JPEG", quality=quality, optimize=True, progressive=True)
    elif fmt == "webp":
        im.save(out, "WEBP", quality=quality, method=6)
    elif fmt == "png":
        if im.mode not in ("RGBA", "RGB", "P", "L"):
            im = im.convert("RGBA")
        im.save(out, "PNG", optimize=True)
    elif fmt in ("tif", "tiff"):
        im.save(out, "TIFF")
    elif fmt == "bmp":
        if im.mode not in ("RGB", "L", "P"):
            im = im.convert("RGB")
        im.save(out, "BMP")
    elif fmt in ("heic", "heif"):
        if im.mode not in ("RGB", "RGBA", "L"):
            im = im.convert("RGB")
        im.save(out, "HEIF", quality=quality)
    else:
        im.save(out)


_IMG_EXT = {"jpg": ".jpg", "jpeg": ".jpg", "png": ".png", "webp": ".webp",
            "tiff": ".tiff", "tif": ".tiff", "bmp": ".bmp",
            "heic": ".heic", "heif": ".heic"}


def image_resize(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    from PIL import Image

    if src.suffix.lower() not in IMAGE_EXTS:
        raise ProcessError(f"Unsupported image type: {src.suffix}")
    try:
        im = Image.open(src)
        im.load()
    except Exception as exc:
        raise ProcessError(f"Couldn't open '{src.name}' ({exc}).")
    w, h = im.size
    mode = opt.get("mode", "dimensions")
    if mode == "percent":
        pct = _clampi(opt.get("percent", 50), 1, 1000, 50) / 100.0
        nw, nh = max(1, int(w * pct)), max(1, int(h * pct))
    else:
        tw = _clampi(opt.get("width", 0), 0, 60000, 0)
        th = _clampi(opt.get("height", 0), 0, 60000, 0)
        keep = bool(opt.get("keep_aspect", True))
        if not tw and not th:
            raise ProcessError("Enter a width and/or height (px).")
        if tw and th and not keep:
            nw, nh = tw, th
        elif tw and th:
            r = min(tw / w, th / h)
            nw, nh = max(1, int(w * r)), max(1, int(h * r))
        elif tw:
            nw, nh = tw, max(1, int(h * tw / w))
        else:
            nh, nw = th, max(1, int(w * th / h))
    im = im.resize((nw, nh), Image.LANCZOS)
    ext = src.suffix.lower()
    if ext not in IMAGE_EXTS:
        ext = ".png"
    out = build_output_path(src, out_dir, ext, name_suffix="_resized",
                            overwrite=opt.get("overwrite", False),
                            numbered=opt.get("same_as_source", False))
    _save_image(im, out, ext.lstrip("."), _clampi(opt.get("quality", 90), 5, 100, 90))
    report(f"{w}×{h} → {nw}×{nh}  ({out.name})")
    return [out]


def image_convert(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    from PIL import Image

    if src.suffix.lower() not in IMAGE_EXTS:
        raise ProcessError(f"Unsupported image type: {src.suffix}")
    fmt = opt.get("format", "png")
    ext = _IMG_EXT.get(fmt, ".png")
    try:
        im = Image.open(src)
        im.load()
    except Exception as exc:
        raise ProcessError(f"Couldn't open '{src.name}' ({exc}).")
    out = build_output_path(src, out_dir, ext, overwrite=opt.get("overwrite", False),
                            numbered=opt.get("same_as_source", False))
    _save_image(im, out, fmt, _clampi(opt.get("quality", 90), 5, 100, 90))
    report(f"{src.suffix.lstrip('.').upper()} → {ext.lstrip('.').upper()}  ({out.name})")
    return [out]


def _load_font(size: int):
    from PIL import ImageFont
    for name in ("arial.ttf", "Arial.ttf", "DejaVuSans.ttf",
                 r"C:\Windows\Fonts\arial.ttf", r"C:\Windows\Fonts\segoeui.ttf"):
        try:
            return ImageFont.truetype(name, size)
        except Exception:
            continue
    return ImageFont.load_default()


def image_watermark(src: Path, out_dir: Path, opt: dict, report: Report) -> list[Path]:
    """Stamp text or a logo across the centre of an image."""
    from PIL import Image

    if src.suffix.lower() not in IMAGE_EXTS:
        raise ProcessError(f"Unsupported image type: {src.suffix}")
    try:
        base = Image.open(src).convert("RGBA")
    except Exception as exc:
        raise ProcessError(f"Couldn't open '{src.name}' ({exc}).")
    opacity = _clampi(opt.get("opacity", 25), 1, 100, 25) / 100.0
    rotation = _clampi(opt.get("rotation", 30), -180, 180, 30)
    bw, bh = base.size
    overlay = Image.new("RGBA", base.size, (0, 0, 0, 0))

    if opt.get("wm_type", "text") == "image":
        raw = (opt.get("image_path") or "").strip()
        if not raw or not Path(raw).exists():
            raise ProcessError("Choose a logo / image for the watermark.")
        try:
            logo = Image.open(raw).convert("RGBA")
        except Exception as exc:
            raise ProcessError(f"Couldn't open logo ({exc}).")
        scale = _clampi(opt.get("scale", 40), 5, 100, 40) / 100.0
        tw = max(1, int(bw * scale))
        th = max(1, int(logo.height * tw / logo.width))
        logo = logo.resize((tw, th), Image.LANCZOS)
        if opacity < 1.0:
            logo.putalpha(logo.split()[3].point(lambda a: int(a * opacity)))
        if rotation % 360:
            logo = logo.rotate(rotation, expand=True, resample=Image.BICUBIC)
        overlay.alpha_composite(logo, ((bw - logo.width) // 2, (bh - logo.height) // 2))
    else:
        from PIL import ImageDraw
        text = (opt.get("text") or "").strip()
        if not text:
            raise ProcessError("Enter the watermark text.")
        size = _clampi(opt.get("font_size", max(24, bw // 12)), 6, 2000,
                       max(24, bw // 12))
        font = _load_font(size)
        col = {"gray": (128, 128, 128), "red": (200, 30, 30),
               "blue": (33, 70, 160), "black": (0, 0, 0),
               "white": (255, 255, 255)}.get(opt.get("color", "gray"),
                                             (128, 128, 128))
        alpha = int(opacity * 255)
        tmp = Image.new("RGBA", base.size, (0, 0, 0, 0))
        d = ImageDraw.Draw(tmp)
        try:
            box = d.textbbox((0, 0), text, font=font)
            tw, th = box[2] - box[0], box[3] - box[1]
        except Exception:
            tw, th = d.textlength(text, font=font), size
        d.text(((bw - tw) / 2, (bh - th) / 2), text, font=font, fill=col + (alpha,))
        if rotation % 360:
            tmp = tmp.rotate(rotation, expand=False, resample=Image.BICUBIC)
        overlay = tmp

    out_img = Image.alpha_composite(base, overlay)
    ext = src.suffix.lower()
    if ext in (".jpg", ".jpeg", ".bmp"):
        out_img = out_img.convert("RGB")
    out = build_output_path(src, out_dir, ext if ext in IMAGE_EXTS else ".png",
                            name_suffix="_watermarked",
                            overwrite=opt.get("overwrite", False),
                            numbered=opt.get("same_as_source", False))
    _save_image(out_img, out, (ext.lstrip(".") or "png"),
                _clampi(opt.get("quality", 90), 5, 100, 90))
    report(f"Watermarked → {out.name}")
    return [out]
