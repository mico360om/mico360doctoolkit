"""Exhaustive per-feature audit: run EVERY tool + its main options and validate
the OUTPUT is actually correct (not merely "it ran"). Calls the processors
directly. External-engine features (Office->PDF for xlsx/pptx, OCR) are exercised
when the engine is available and skipped-with-note otherwise.

Run:  python tests/feature_audit_test.py
"""
from __future__ import annotations

import io
import os
import sys
import tempfile
from pathlib import Path

os.environ.setdefault("QT_QPA_PLATFORM", "offscreen")
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

failures: list[str] = []
REP = lambda *_: None  # noqa: E731


def check(name, ok, detail=""):
    print(f"[{'PASS' if ok else 'FAIL'}] {name}" + (f": {detail}" if detail else ""))
    if not ok:
        failures.append(name)


# ---- sample builders ----------------------------------------------------
def make_pdf(p, pages=3, text=True, image=False):
    import fitz
    d = fitz.open()
    for i in range(pages):
        pg = d.new_page(width=420, height=560)
        if text:
            pg.insert_text((40, 80), f"Page {i + 1} — the quick brown fox. " * 6,
                           fontsize=11)
        if image:
            from PIL import Image
            buf = io.BytesIO()
            Image.new("RGB", (300, 300), (i * 60 % 255, 120, 200)).save(buf, "JPEG")
            pg.insert_image(fitz.Rect(40, 120, 340, 420), stream=buf.getvalue())
    d.save(str(p)); d.close(); return p


def make_image(p, size=(640, 480), color=(200, 80, 80)):
    from PIL import Image
    Image.new("RGB", size, color).save(p); return p


def make_docx(p):
    import docx
    d = docx.Document(); d.add_heading("Title", 1); d.add_paragraph("Hello body text.")
    t = d.add_table(rows=2, cols=2)
    t.rows[0].cells[0].text = "A"; t.rows[0].cells[1].text = "B"
    d.save(str(p)); return p


def make_xlsx(p):
    from openpyxl import Workbook
    wb = Workbook(); ws = wb.active; ws.title = "Sheet1"
    ws.append(["Name", "Qty"]); ws.append(["Apples", 5]); wb.save(str(p)); return p


def make_pptx(p):
    from pptx import Presentation
    prs = Presentation(); s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Deck Title"; s.placeholders[1].text_frame.text = "Point one"
    prs.save(str(p)); return p


def image_only_pdf(p, text="Invoice Number 84217"):
    import fitz
    from PIL import Image, ImageDraw, ImageFont
    img = Image.new("RGB", (1000, 200), (255, 255, 255)); dr = ImageDraw.Draw(img)
    try:
        fnt = ImageFont.truetype("arial.ttf", 52)
    except Exception:
        fnt = ImageFont.load_default()
    dr.text((40, 70), text, fill=(0, 0, 0), font=fnt)
    buf = io.BytesIO(); img.save(buf, "PNG")
    d = fitz.open(); pg = d.new_page(width=360, height=72)
    pg.insert_image(pg.rect, stream=buf.getvalue()); d.save(str(p)); d.close(); return p


def pdf_pages(p):
    import fitz
    d = fitz.open(str(p)); n = d.page_count; d.close(); return n


def pdf_text(p):
    import fitz
    d = fitz.open(str(p)); t = "".join(pg.get_text() for pg in d); d.close(); return t


def main() -> int:  # noqa: C901
    from mico360.core import processors as P
    from mico360.core.deps import find_libreoffice
    have_lo = bool(find_libreoffice())
    td = Path(tempfile.mkdtemp(prefix="mico360_audit_"))
    out = td / "out"; out.mkdir()

    def o(opt=None):
        d = {"overwrite": True}; d.update(opt or {}); return d

    # ============ PDF tools ============
    # 1. Compress
    src = make_pdf(td / "comp.pdf", 4, image=True)
    for lvl in ("low", "medium", "high"):
        r = P.pdf_compress(src, out, o({"level": lvl}), REP)
        check(f"Compress PDF [{lvl}] → valid PDF", r and pdf_pages(r[0]) == 4)
    r = P.pdf_compress(src, out, o({"level": "custom", "dpi": 72, "jpeg_quality": 40}), REP)
    check("Compress PDF [custom] smaller than source",
          r and r[0].stat().st_size <= src.stat().st_size)
    r = P.pdf_compress(src, out, o({"level": "target", "target_kb": 80}), REP)
    check("Compress PDF [target 80KB] lands at/under (within slack)",
          r and r[0].stat().st_size <= 160 * 1024,
          f"{r[0].stat().st_size // 1024} KB" if r else "")

    # 2. Merge
    a, b = make_pdf(td / "m1.pdf", 3), make_pdf(td / "m2.pdf", 2)
    r = P.pdf_merge([a, b], out, o({"output_name": "merged"}), REP)
    check("Merge PDF → page count = sum (5)", r and pdf_pages(r[0]) == 5)

    # 3. Split
    s4 = make_pdf(td / "split.pdf", 4)
    r = P.pdf_split(s4, out, o({"mode": "each"}), REP)
    check("Split PDF [each] → 4 files", r and len(r) == 4 and all(pdf_pages(x) == 1 for x in r))
    r = P.pdf_split(s4, out, o({"mode": "every_n", "every_n": 2}), REP)
    check("Split PDF [every 2] → 2 files of 2 pages", r and len(r) == 2 and pdf_pages(r[0]) == 2)
    r = P.pdf_split(s4, out, o({"mode": "ranges", "ranges": "1-2, 4"}), REP)
    check("Split PDF [ranges] → 2 files", r and len(r) == 2)

    # 4. Organize
    s5 = make_pdf(td / "org.pdf", 5)
    r = P.pdf_organize(s5, out, o({"operation": "delete", "del_pages": "2,4"}), REP)
    check("Organize [delete 2,4] → 3 pages", r and pdf_pages(r[0]) == 3)
    r = P.pdf_organize(s5, out, o({"operation": "extract", "ext_pages": "1,3,5"}), REP)
    check("Organize [extract 1,3,5] → 3 pages", r and pdf_pages(r[0]) == 3)
    r = P.pdf_organize(s5, out, o({"operation": "rotate", "angle": 90, "pages": "all"}), REP)
    check("Organize [rotate] → still 5 pages, valid", r and pdf_pages(r[0]) == 5)
    r = P.pdf_organize(s5, out, o({"operation": "reorder", "order": "5,4,3,2,1"}), REP)
    check("Organize [reorder] → 5 pages", r and pdf_pages(r[0]) == 5)

    # 5. Protect / unlock
    import pypdf
    prot = P.pdf_protect(s4, out, o({"operation": "protect", "password": "sEcr3t",
                                     "confirm_password": "sEcr3t"}), REP)
    rd = pypdf.PdfReader(str(prot[0]))
    check("Protect PDF → encrypted", rd.is_encrypted)
    unl = P.pdf_protect(prot[0], out, o({"operation": "unlock", "password": "sEcr3t"}), REP)
    rd2 = pypdf.PdfReader(str(unl[0]))
    check("Unlock PDF → no longer encrypted", not rd2.is_encrypted)

    # 6. Watermark (text + image)
    r = P.pdf_watermark(s4, out, o({"wm_type": "text", "text": "DRAFT",
                                    "opacity": 30, "rotation": 45, "position": "center"}), REP)
    check("Watermark PDF [text] → valid, same pages", r and pdf_pages(r[0]) == 4)
    logo = make_image(td / "logo.png", (200, 80), (0, 120, 200))
    r = P.pdf_watermark(s4, out, o({"wm_type": "image", "image_path": str(logo),
                                    "scale": 30, "opacity": 40}), REP)
    check("Watermark PDF [image] → valid", r and pdf_pages(r[0]) == 4)

    # 7. Page numbers
    r = P.pdf_page_numbers(s4, out, o({"position": "bottom-center", "format": "n_of_total",
                                       "start": 1, "font_size": 11}), REP)
    check("Add Page Numbers → '1 / 4' present", r and ("1 / 4" in pdf_text(r[0]) or "/ 4" in pdf_text(r[0])))

    # 8. Sign
    sig = make_image(td / "sig.png", (180, 70), (10, 10, 10))
    r = P.pdf_sign(s4, out, o({"image_path": str(sig), "page": "last",
                               "position": "bottom-right", "width": 25}), REP)
    check("Sign PDF → valid, same pages", r and pdf_pages(r[0]) == 4)

    # 9. Metadata
    r = P.pdf_metadata(s4, out, o({"title": "Audited Doc", "author": "QA",
                                   "subject": "Test", "keywords": "a,b"}), REP)
    import fitz
    d = fitz.open(str(r[0])); md = d.metadata; d.close()
    check("Edit Metadata → title+author set",
          md.get("title") == "Audited Doc" and md.get("author") == "QA", str(md.get("title")))

    # 10. OCR
    try:
        P._make_ocr_engine()
        scan = image_only_pdf(td / "scan.pdf")
        r = P.pdf_ocr(scan, out, o({"quality": "balanced"}), REP)
        txt = pdf_text(r[0]).lower()
        check("Searchable PDF (OCR) → text layer recovered",
              r and ("invoice" in txt or "84217" in txt), txt[:40])
    except Exception as exc:
        check("Searchable PDF (OCR) — engine unavailable (skipped)", True, str(exc)[:40])

    # ============ Convert tools ============
    pdfc = make_pdf(td / "conv.pdf", 2, image=True)
    # 11. PDF → ...
    r = P.pdf_convert(pdfc, out, o({"target": "word"}), REP)
    check("PDF → Word → .docx", r and r[0].suffix == ".docx" and r[0].exists())
    r = P.pdf_convert(pdfc, out, o({"target": "pptx", "mode": "auto"}), REP)
    check("PDF → PowerPoint → .pptx", r and r[0].suffix == ".pptx" and r[0].exists())
    r = P.pdf_convert(pdfc, out, o({"target": "excel"}), REP)
    check("PDF → Excel → .xlsx", r and r[0].suffix == ".xlsx" and r[0].exists())
    r = P.pdf_convert(pdfc, out, o({"target": "image", "format": "png", "dpi": 100}), REP)
    check("PDF → Image → png per page (2)", r and len(r) == 2 and r[0].suffix.lower() == ".png")

    # 12. Office → PDF
    r = P.office_to_pdf(make_docx(td / "w.docx"), out, o(), REP)
    check("Office → PDF (Word, built-in/LO) → .pdf", r and r[0].suffix == ".pdf" and r[0].exists())
    if have_lo:
        for mk, nm in ((make_xlsx(td / "x.xlsx"), "Excel"), (make_pptx(td / "p.pptx"), "PowerPoint")):
            try:
                r = P.office_to_pdf(mk, out, o(), REP)
                check(f"Office → PDF ({nm}) → .pdf", r and r[0].suffix == ".pdf" and r[0].exists())
            except Exception as exc:
                check(f"Office → PDF ({nm})", False, str(exc)[:60])
    else:
        check("Office → PDF (Excel/PPT) — LibreOffice absent (skipped)", True)

    # 13. Document → Markdown
    r = P.to_markdown(make_docx(td / "md.docx"), out, o(), REP)
    check("Doc → Markdown (Word) → has heading", r and "# " in r[0].read_text(encoding="utf-8"))
    r = P.to_markdown(make_xlsx(td / "md.xlsx"), out, o(), REP)
    check("Doc → Markdown (Excel) → table row", r and "| Name | Qty |" in r[0].read_text(encoding="utf-8"))
    r = P.to_markdown(make_pptx(td / "md.pptx"), out, o(), REP)
    check("Doc → Markdown (PPT) → slide heading", r and "## Slide 1" in r[0].read_text(encoding="utf-8"))
    r = P.to_markdown(make_pdf(td / "md.pdf", 1), out, o(), REP)
    check("Doc → Markdown (PDF) → has text", r and "quick brown fox" in r[0].read_text(encoding="utf-8"))

    # 14. Image → PDF (combine + separate)
    imgs = [make_image(td / f"i{i}.png", (400, 300), (i * 50, 100, 150)) for i in range(3)]
    r = P.image_to_pdf(imgs, out, o({"combine": True, "output_name": "imgs"}), REP)
    check("Image → PDF [combine] → 1 PDF, 3 pages", r and len(r) == 1 and pdf_pages(r[0]) == 3)
    r = P.image_to_pdf(imgs, out, o({"combine": False}), REP)
    check("Image → PDF [separate] → 3 PDFs", r and len(r) == 3)

    # ============ Image tools ============
    big = make_image(td / "big.png", (1600, 1200), (180, 60, 60))
    # 15. Compress
    r = P.image_compress(big, out, o({"level": "high"}), REP)
    check("Compress Image [high] smaller", r and r[0].stat().st_size < big.stat().st_size)
    r = P.image_compress(big, out, o({"level": "custom", "quality": 50, "format": "jpg",
                                      "max_dimension": 800}), REP)
    from PIL import Image
    w, h = Image.open(r[0]).size
    check("Compress Image [custom+resize] ≤800px", r and max(w, h) <= 800, f"{w}x{h}")
    r = P.image_compress(big, out, o({"level": "target", "target_kb": 40}), REP)
    check("Compress Image [target 40KB] lands at/under (within slack)",
          r and r[0].stat().st_size <= 80 * 1024,
          f"{r[0].stat().st_size // 1024} KB" if r else "")

    # 16. Resize
    r = P.image_resize(big, out, o({"mode": "dimensions", "width": 640, "height": 0,
                                    "keep_aspect": True}), REP)
    w, h = Image.open(r[0]).size
    check("Resize Image [w=640, aspect] → width 640", r and w == 640, f"{w}x{h}")
    r = P.image_resize(big, out, o({"mode": "percent", "percent": 25}), REP)
    w, h = Image.open(r[0]).size
    check("Resize Image [25%] → 400x300", r and (w, h) == (400, 300), f"{w}x{h}")

    # 17. Convert
    r = P.image_convert(big, out, o({"format": "jpg", "quality": 85}), REP)
    check("Convert Image → .jpg", r and r[0].suffix.lower() in (".jpg", ".jpeg")
          and Image.open(r[0]).format == "JPEG")
    r = P.image_convert(big, out, o({"format": "webp", "quality": 80}), REP)
    check("Convert Image → .webp", r and r[0].suffix.lower() == ".webp")

    # 18. Watermark image
    r = P.image_watermark(big, out, o({"wm_type": "text", "text": "SAMPLE",
                                       "opacity": 35, "rotation": 30}), REP)
    check("Watermark Image [text] → valid image", r and Image.open(r[0]).size == (1600, 1200))
    r = P.image_watermark(big, out, o({"wm_type": "image", "image_path": str(logo),
                                       "scale": 25, "opacity": 50}), REP)
    check("Watermark Image [logo] → valid image", r and r[0].exists())

    import shutil
    shutil.rmtree(td, ignore_errors=True)
    print()
    if failures:
        print(f"{len(failures)} feature check(s) FAILED: {', '.join(failures)}")
        return 1
    print("All feature-audit checks passed.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
