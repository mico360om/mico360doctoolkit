"""AI metadata generator + user AI configuration.

Runs entirely offline: a stub HTTP server stands in for the AI provider, so the
whole path (config -> request -> parse -> suggestion panel -> apply) is exercised
without a key or a network. Key storage is asserted to never keep clear text.

Run:  python tests/ai_metadata_test.py
"""
from __future__ import annotations

import json
import os
import sys
import tempfile
import threading
from http.server import BaseHTTPRequestHandler, HTTPServer
from pathlib import Path

os.environ.setdefault("QT_QPA_PLATFORM", "offscreen")
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

failures: list[str] = []


def check(name, ok, detail=""):
    print(f"[{'PASS' if ok else 'FAIL'}] {name}" + (f": {detail}" if detail else ""))
    if not ok:
        failures.append(name)


GOOD_KEY = "mico_test_key_abcdefghijklmnop"
REPLY = {
    "title": "Quarterly Portfolio Review",
    "author": "Strategy Office",
    "subject": "Performance of the 2026 portfolio",
    "keywords": "portfolio, governance, review",
    "creator": "Microsoft Word",
    "producer": "Acrobat Distiller",
    "creation_date": "2026-03-04",
    "mod_date": "2026-03-09",
    "company": "MICO360",
    "manager": "A. Manager",
    "category": "Report",
    "comments": "A quarterly review of portfolio performance and risks.",
    "custom": "Invoice Number = INV-1042",
    "copyright": "(c) 2026 MICO360",
    "language": "en-US",
    "trapped": "Unknown",
    # Values the AI is NOT confident about must come back blank and must never
    # be applied over an existing value.
    "_unused": "",
}


class Handler(BaseHTTPRequestHandler):
    def log_message(self, *a):
        pass

    def _json(self, code, obj):
        body = json.dumps(obj).encode()
        self.send_response(code)
        self.send_header("Content-Type", "application/json")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    def _authed(self) -> bool:
        return self.headers.get("Authorization") == f"Bearer {GOOD_KEY}"

    def do_GET(self):
        if not self._authed():
            return self._json(401, {"error": {"message": "Invalid or missing API key."}})
        if self.path.rstrip("/").endswith("/v1/models"):
            return self._json(200, {"object": "list", "data": [
                {"id": "qwen2.5:0.5b", "object": "model"},
                {"id": "qwen2.5vl:7b", "object": "model"},
                {"id": "llama3.1:8b", "object": "model"}]})
        self._json(404, {"error": {"message": "no route"}})

    def do_POST(self):
        if not self._authed():
            return self._json(401, {"error": {"message": "Invalid or missing API key."}})
        if not self.path.rstrip("/").endswith("/v1/chat/completions"):
            return self._json(404, {"error": {"message": "no route"}})
        n = int(self.headers.get("Content-Length") or 0)
        body = json.loads(self.rfile.read(n) or b"{}")
        # The document excerpt must actually reach the model.
        sent = " ".join(m.get("content", "") for m in body.get("messages", [])
                        if isinstance(m.get("content"), str))
        content = ("Here you go:\n```json\n" + json.dumps(REPLY) + "\n```"
                   if "PORTFOLIO REVIEW" in sent.upper() else "{}")
        self._json(200, {"choices": [{"index": 0, "message":
                                      {"role": "assistant", "content": content},
                                      "finish_reason": "stop"}]})


def main() -> int:
    from PySide6.QtWidgets import QApplication
    app = QApplication.instance() or QApplication([])

    from mico360.config import settings
    from mico360.core import ai as ai_core
    from mico360.core import ai_metadata

    srv = HTTPServer(("127.0.0.1", 0), Handler)
    threading.Thread(target=srv.serve_forever, daemon=True).start()
    base = f"http://127.0.0.1:{srv.server_address[1]}/v1"

    saved = (settings.ai_enabled, settings.ai_source, settings.ai_base_url,
             settings.ai_model, settings.ai_api_key_sealed)
    try:
        # ---------- key security ------------------------------------
        sealed = ai_core.seal_key(GOOD_KEY)
        check("stored key is never clear text", GOOD_KEY not in sealed, sealed[:12])
        check("stored key decrypts back correctly",
              ai_core.unseal_key(sealed) == GOOD_KEY)
        check("masked key hides the secret",
              GOOD_KEY not in ai_core.masked_key(GOOD_KEY)
              and ai_core.masked_key(GOOD_KEY).endswith(GOOD_KEY[-4:]),
              ai_core.masked_key(GOOD_KEY))
        check("a corrupt stored key fails safely (no crash)",
              ai_core.unseal_key("dpapi:not-base64!!") == "")

        # ---------- URL normalisation --------------------------------
        cases = {
            "ai.example.com:5310": "http://ai.example.com:5310/v1",
            "http://x.test:5310/": "http://x.test:5310/v1",
            "http://x.test:5310/v1/": "http://x.test:5310/v1",
            "http://x.test:5310/v1/chat/completions": "http://x.test:5310/v1",
        }
        bad = {k: ai_core.normalize_base_url(k) for k, v in cases.items()
               if ai_core.normalize_base_url(k) != v}
        check("base URLs are normalised (port + /v1, no duplicate path)",
              not bad, str(bad))

        # ---------- config gating ------------------------------------
        cfg = ai_core.AiConfig(enabled=False)
        check("disabled AI reports why", cfg.is_usable()[0] is False
              and "turned off" in cfg.is_usable()[1])
        cfg = ai_core.AiConfig(enabled=True, source=ai_core.SOURCE_CUSTOM,
                               base_url=base, api_key="")
        check("missing key reports 'not configured'",
              cfg.is_usable()[0] is False and "key" in cfg.is_usable()[1].lower())

        # ---------- connection test ----------------------------------
        good = ai_core.AiConfig(enabled=True, source=ai_core.SOURCE_CUSTOM,
                                base_url=base, api_key=GOOD_KEY,
                                model="qwen2.5:0.5b")
        ok, msg = ai_core.test_connection(good)
        check("Test connection succeeds and names the model", ok and "qwen2.5" in msg,
              msg)
        check("the server's full model list is returned",
              set(ai_core.list_models(good)) ==
              {"qwen2.5:0.5b", "qwen2.5vl:7b", "llama3.1:8b"})

        wrong = ai_core.AiConfig(enabled=True, source=ai_core.SOURCE_CUSTOM,
                                 base_url=base, api_key="nope")
        ok, msg = ai_core.test_connection(wrong)
        check("a rejected key gives actionable advice",
              not ok and "key" in msg.lower(), msg[:60])

        dead = ai_core.AiConfig(enabled=True, source=ai_core.SOURCE_CUSTOM,
                                base_url="http://127.0.0.1:9/v1", api_key="x")
        ok, msg = ai_core.test_connection(dead)
        check("an unreachable server fails gracefully",
              not ok and "couldn't reach" in msg.lower(), msg[:60])

        # ---------- extraction + suggestion --------------------------
        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            import fitz
            doc = fitz.open()
            page = doc.new_page()
            page.insert_text((60, 90), "PORTFOLIO REVIEW", fontsize=18)
            page.insert_text((60, 130),
                             "This document reviews the 2026 portfolio "
                             "performance across all programmes.", fontsize=11)
            pdf = tmp / "review.pdf"
            doc.save(str(pdf))
            doc.close()

            text = ai_metadata.extract_text(pdf)
            check("document text is extracted locally",
                  "PORTFOLIO REVIEW" in text.upper(), text[:40])

            scan = tmp / "scan.pdf"
            d2 = fitz.open()
            d2.new_page()
            d2.save(str(scan))
            d2.close()
            try:
                ai_metadata.extract_text(scan)
                check("a text-less PDF advises running OCR", False, "no error")
            except ai_metadata.NoTextError as exc:
                check("a text-less PDF advises running OCR",
                      "searchable pdf" in str(exc).lower(), str(exc)[:60])

            # An image-only Office file must get advice that FITS ITS TYPE —
            # the OCR tool only accepts PDFs, so "run OCR on it" would be
            # useless for a .pptx.
            from pptx import Presentation
            deck = tmp / "images_only.pptx"
            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[6])   # blank, no text
            prs.save(str(deck))
            try:
                ai_metadata.extract_text(deck)
                check("an image-only deck is refused with usable advice",
                      False, "no error")
            except ai_metadata.NoTextError as exc:
                m = str(exc)
                check("an image-only deck is refused with usable advice",
                      "Office" in m and "PDF" in m, m[:80])

            got = ai_metadata.suggest_metadata(pdf, good)
            check("AI returns every supported metadata field",
                  set(got) == set(ai_metadata.FIELDS),
                  f"{len(got)}/{len(ai_metadata.FIELDS)}: {sorted(got)}")
            check("dates come back normalised", got["creation_date"] == "2026-03-04")
            check("trapped is one of the allowed values",
                  got["trapped"] in ("True", "False", "Unknown"))
            check("custom properties keep their Key = Value form",
                  "=" in got["custom"], got["custom"])
            check("suggested title matches the document",
                  got["title"] == REPLY["title"], got.get("title"))
            check("JSON wrapped in prose/code fences is still parsed",
                  got["language"] == "en-US")

            # ---------- panel: review, edit, apply -------------------
            settings.ai_enabled = True
            settings.ai_source = ai_core.SOURCE_CUSTOM
            settings.ai_base_url = base
            settings.ai_model = "qwen2.5:0.5b"
            settings.ai_api_key_sealed = ai_core.seal_key(GOOD_KEY)

            from mico360.core.tools import TOOLS_BY_ID
            from mico360.ui.tool_page import ToolPage
            page_ui = ToolPage(TOOLS_BY_ID["pdf_metadata"])
            check("Edit Metadata page shows the AI panel",
                  page_ui.ai_panel is not None)
            panel = page_ui.ai_panel
            panel.refresh_availability()
            check("panel reports ready when AI is configured",
                  panel.btn_suggest.isEnabled()
                  and "not configured" not in panel.status.text().lower(),
                  panel.status.text()[:60])

            # Suggestions are shown, NOT auto-applied. (Option fields may hold
            # remembered values, so compare against what was there before.)
            title_ctrl = page_ui.options_widget._controls["title"]
            before_title = title_ctrl.text()
            panel._show(got)
            check("suggestions are shown for review, not applied automatically",
                  title_ctrl.text() == before_title, repr(title_ctrl.text()))
            # NOTE: isVisible() is False for any widget whose window isn't shown,
            # so assert on the not-hidden state instead.
            check("every suggested field gets a row",
                  not panel.results.isHidden()
                  and all(panel._rows[k][2].text() for k in got))

            # Apply one field.
            author_ctrl = page_ui.options_widget._controls["author"]
            before_author = author_ctrl.text()
            panel._apply_one("title")
            check("applying one field fills ONLY that option",
                  title_ctrl.text() == REPLY["title"]
                  and author_ctrl.text() == before_author,
                  f"title={title_ctrl.text()!r} author={author_ctrl.text()!r}")

            # Edit a suggestion, then apply all.
            panel._rows["author"][2].setText("Edited Author")
            panel._apply_all()
            vals = page_ui.options_widget.values()
            check("Apply all uses the EDITED value",
                  vals.get("author") == "Edited Author", vals.get("author"))
            check("Apply all fills the remaining fields",
                  vals.get("keywords") == REPLY["keywords"]
                  and vals.get("category") == REPLY["category"])
            check("suggested comments land in the custom Comments field",
                  vals.get("comments") == REPLY["comments"])

            # Dismiss hides them.
            panel.clear()
            check("Dismiss hides the suggestions and forgets them",
                  panel.results.isHidden() and not panel._current)

            # ---------- not configured -> clear message + link -------
            settings.ai_enabled = False
            panel.refresh_availability()
            check("turning AI off shows 'AI API not configured'",
                  "not configured" in panel.status.text().lower()
                  and not panel.btn_suggest.isEnabled(), panel.status.text()[:50])
            check("a Configure button is offered when unconfigured",
                  not panel.btn_configure.isHidden())

            # ---------- settings page round-trip ---------------------
            settings.ai_enabled = True
            from mico360.ui.settings_page import SettingsPage
            sp = SettingsPage()
            check("Settings has the AI fields",
                  all(hasattr(sp, a) for a in
                      ("chk_ai", "ai_source", "ai_url", "ai_key", "ai_model",
                       "btn_ai_test")))
            check("the saved key is never rendered in full",
                  GOOD_KEY not in sp.ai_key.text()
                  and GOOD_KEY not in sp.ai_key.placeholderText(),
                  sp.ai_key.placeholderText())
            check("the key field is masked input",
                  sp.ai_key.echoMode() == type(sp.ai_key).Password)
            sp.ai_key.setText("mico_a_brand_new_key_value")
            sp._save_ai()
            check("saving a new key encrypts it and clears the box",
                  sp.ai_key.text() == ""
                  and ai_core.unseal_key(settings.ai_api_key_sealed)
                  == "mico_a_brand_new_key_value")
            check("leaving the key blank keeps the existing one",
                  (sp._save_ai() or True)
                  and ai_core.unseal_key(settings.ai_api_key_sealed)
                  == "mico_a_brand_new_key_value")

            # --- Test connection: BOTH outcomes must render, not crash ------
            # (A missing palette key once made the FAILURE path raise KeyError,
            #  which is exactly the path a misconfigured user hits.)
            for theme in ("light", "dark"):
                prev_theme = settings.theme_mode
                settings.theme_mode = theme
                sp.ai_source.setCurrentIndex(
                    sp.ai_source.findData(ai_core.SOURCE_CUSTOM))
                # success
                sp.ai_url.setText(base)
                sp.ai_key.setText(GOOD_KEY)
                sp._test_ai()
                ok_txt = sp.ai_status.text()
                check(f"Test connection renders SUCCESS ({theme} theme)",
                      "qwen2.5" in ok_txt and "<span" in ok_txt, ok_txt[:60])
                # failure — unreachable server
                sp.ai_url.setText("http://127.0.0.1:9/v1")
                sp._test_ai()
                bad_txt = sp.ai_status.text()
                check(f"Test connection renders FAILURE without crashing ({theme})",
                      "reach" in bad_txt.lower() and "<span" in bad_txt,
                      bad_txt[:70])
                settings.theme_mode = prev_theme

            # A server message containing markup must not corrupt the label.
            sp.ai_status.setText("")
            import mico360.core.ai as _ai
            _orig = _ai.test_connection
            try:
                _ai.test_connection = lambda cfg: (False, "bad <b>url</b> & key")
                sp._test_ai()
                check("a message with markup is escaped, not rendered",
                      "&lt;b&gt;" in sp.ai_status.text(), sp.ai_status.text()[:70])
            finally:
                _ai.test_connection = _orig

            # --- model dropdown: list / add / remove --------------------
            from PySide6.QtWidgets import QComboBox
            check("Model is a dropdown you can also type into",
                  isinstance(sp.ai_model, QComboBox) and sp.ai_model.isEditable())

            sp.ai_source.setCurrentIndex(
                sp.ai_source.findData(ai_core.SOURCE_CUSTOM))
            sp.ai_url.setText(base)
            sp.ai_key.setText(GOOD_KEY)
            sp._refresh_models()
            listed = [sp.ai_model.itemText(i) for i in range(sp.ai_model.count())]
            check("Refresh lists every model the server offers",
                  {"qwen2.5:0.5b", "qwen2.5vl:7b", "llama3.1:8b"} <= set(listed),
                  str(listed))
            check("the model list is remembered in settings",
                  "llama3.1:8b" in settings.ai_models)

            # Add: type an id the server doesn't know, save -> it joins the list.
            sp.ai_model.setEditText("my-private-model:latest")
            sp._save_ai()
            check("typing a new model id adds it to the list",
                  "my-private-model:latest" in settings.ai_models
                  and settings.ai_model == "my-private-model:latest")
            sp._load_models()
            check("the added model appears in the dropdown",
                  "my-private-model:latest" in
                  [sp.ai_model.itemText(i) for i in range(sp.ai_model.count())])

            # Remove: drops it from the dropdown, not from the server.
            i = sp.ai_model.findText("my-private-model:latest")
            sp.ai_model.setCurrentIndex(i)
            sp._remove_model()
            after = [sp.ai_model.itemText(i) for i in range(sp.ai_model.count())]
            check("Remove takes the model out of the list",
                  "my-private-model:latest" not in after
                  and "my-private-model:latest" not in settings.ai_models,
                  str(after))
            check("removing never empties the dropdown", len(after) >= 1)

            # A refresh against an unreachable server must not crash.
            sp.ai_url.setText("http://127.0.0.1:9/v1")
            sp._refresh_models()
            check("Refresh against a dead server reports it, no crash",
                  "couldn" in sp.ai_models_state.text().lower(),
                  sp.ai_models_state.text()[:60])

            # ===== panel: all fields, bulk update, auto apply ==========
            panel = page_ui.ai_panel
            page_ui.options_widget._controls["privacy"].setCurrentIndex(0)
            check("a row exists for every metadata field",
                  len(panel._rows) == len(ai_metadata.FIELDS))
            check("the summary counts this tool's AI-fillable fields",
                  panel._total_fields == len(ai_metadata.FIELDS),
                  str(panel._total_fields))

            panel._show(got)
            check("all suggested fields are shown", len(panel._current) == len(got))

            # Bulk update: tick only two fields, apply just those.
            for k in panel._current:
                panel._rows[k][0].setChecked(False)
            panel._rows["company"][0].setChecked(True)
            panel._rows["manager"][0].setChecked(True)
            for k in ("company", "manager", "copyright"):
                page_ui.options_widget._controls[k].setText("")
            panel._apply_selected()
            vals = page_ui.options_widget.values()
            check("Apply selected updates ONLY the ticked fields",
                  vals["company"] == "MICO360" and vals["manager"] == "A. Manager"
                  and vals["copyright"] == "",
                  f"copyright={vals['copyright']!r}")
            check("the summary reports how many fields were updated",
                  "2 of" in panel.status.text(), panel.status.text()[:60])

            # Select all -> apply all covers the rest, including the choice and
            # multi-line controls.
            panel._toggle_select_all()
            panel._apply_all()
            vals = page_ui.options_widget.values()
            check("Apply all fills every suggested field",
                  vals["copyright"] == "(c) 2026 MICO360"
                  and vals["creation_date"] == "2026-03-04")
            check("a choice field (Trapped) takes the suggestion",
                  vals["trapped"] == "Unknown", str(vals["trapped"]))
            check("a multi-line field (Custom) takes the suggestion",
                  "Invoice Number = INV-1042" in vals["custom"], vals["custom"])

            # Never overwrite a good value with a blank suggestion.
            page_ui.options_widget._controls["title"].setText("Keep me")
            page_ui._apply_ai_field("title", "   ")
            check("a blank suggestion never wipes an existing value",
                  page_ui.options_widget.values()["title"] == "Keep me")

            # Privacy preset must be respected.
            pv = page_ui.options_widget._controls["privacy"]
            pv.setCurrentIndex(pv.findData("scrub"))
            page_ui.options_widget._controls["author"].setText("Original")
            page_ui._apply_ai_field("author", "AI Person")
            check("suggestions are not applied while Privacy is on",
                  page_ui.options_widget.values()["author"] == "Original")
            pv.setCurrentIndex(pv.findData(""))

            # Auto apply: suggestions land without any confirmation.
            for k in ("company", "category"):
                page_ui.options_widget._controls[k].setText("")
            panel.chk_auto.setChecked(True)
            panel._on_done(dict(got))
            vals = page_ui.options_widget.values()
            check("Auto apply fills the fields with no confirmation",
                  vals["company"] == "MICO360" and vals["category"] == "Report")
            check("auto-apply preference is remembered", settings.ai_auto_apply)
            panel.chk_auto.setChecked(False)
            check("turning auto-apply off is remembered too",
                  not settings.ai_auto_apply)
    finally:
        srv.shutdown()
        (settings.ai_enabled, settings.ai_source, settings.ai_base_url,
         settings.ai_model, settings.ai_api_key_sealed) = saved

    print()
    if failures:
        print(f"{len(failures)} AI check(s) FAILED: {', '.join(failures)}")
        return 1
    print("All AI metadata / configuration checks passed.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
