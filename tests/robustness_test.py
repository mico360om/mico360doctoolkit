"""Import & processing robustness: the app must never crash, hang, or lose work
because of a bad file.

For EVERY registered tool this queues a mix of valid + corrupt + empty +
truncated files and asserts that the batch completes, the good file still
succeeds, and every bad file is flagged as failed with a readable message —
i.e. problems are skipped/flagged and the rest keeps processing.

Run:  python tests/robustness_test.py
"""
from __future__ import annotations

import os
import sys
import tempfile
from pathlib import Path

os.environ.setdefault("QT_QPA_PLATFORM", "offscreen")
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace", line_buffering=True)
except Exception:
    pass

from PySide6.QtCore import QEventLoop, QTimer                       # noqa: E402
from PySide6.QtWidgets import (QApplication, QCheckBox, QComboBox,   # noqa: E402
                               QLineEdit, QMessageBox, QSpinBox)

from mico360.config import settings                                 # noqa: E402
from mico360.core.tools import (AGGREGATE, EXCEL, IMAGES, PDF, PPT,  # noqa: E402
                                TOOLS, WORD)
from mico360.ui.tool_page import ToolPage                           # noqa: E402

failures: list[str] = []


def check(name: str, ok: bool, detail: str = "") -> None:
    print(f"[{'PASS' if ok else 'FAIL'}] {name}" + (f": {detail}" if detail else ""))
    if not ok:
        failures.append(name)


def _set(page, key, value):
    ctrl = page.options_widget._controls.get(key)
    if ctrl is None:
        return
    if isinstance(ctrl, QComboBox):
        i = ctrl.findData(value)
        if i >= 0:
            ctrl.setCurrentIndex(i)
    elif isinstance(ctrl, QSpinBox):
        ctrl.setValue(int(value))
    elif isinstance(ctrl, QCheckBox):
        ctrl.setChecked(bool(value))
    elif isinstance(ctrl, QLineEdit):
        ctrl.setText(str(value))


def run(page, out, timeout_ms=90000):
    page.chk_same.setChecked(False)
    page.chk_overwrite.setChecked(True)
    page.out_edit.setText(str(out))
    res = {}
    loop = QEventLoop()
    page.start()
    if page.controller is None:
        return {"_nostart": True}
    page.controller.finished.connect(lambda s: (res.update(s), loop.quit()))
    g = QTimer()
    g.setSingleShot(True)
    g.timeout.connect(loop.quit)
    g.start(timeout_ms)
    loop.exec()
    res["_timeout"] = not g.isActive()
    g.stop()
    return res


# --- good samples ---------------------------------------------------------
def make_pdf(p, pages=2):
    import fitz
    d = fitz.open()
    for i in range(pages):
        d.new_page(width=400, height=300).insert_text((50, 80), f"PAGE {i+1}",
                                                      fontsize=18)
    d.save(str(p))
    d.close()
    return p


def make_png(p):
    from PIL import Image
    Image.new("RGB", (400, 300), (90, 140, 200)).save(p)
    return p


def make_docx(p):
    import docx
    doc = docx.Document()
    doc.add_paragraph("Hello")
    doc.save(str(p))
    return p


def make_xlsx(p):
    from openpyxl import Workbook
    wb = Workbook()
    wb.active.append(["A", 1])
    wb.save(str(p))
    return p


def make_pptx(p):
    from pptx import Presentation
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(p))
    return p


def make_svg(p):
    p.write_text('<svg xmlns="http://www.w3.org/2000/svg" width="80" height="60">'
                 '<rect width="80" height="60" fill="#A0201F"/></svg>')
    return p


_MAKERS = {".svg": make_svg, ".png": make_png, ".pdf": make_pdf,
           ".docx": make_docx, ".xlsx": make_xlsx, ".pptx": make_pptx}


def ext_for(tool) -> str:
    a = tool.accept
    if ".svg" in a:
        return ".svg"
    if a & IMAGES:
        return ".png"
    if a & PDF:
        return ".pdf"
    if a & WORD:
        return ".docx"
    if a & EXCEL:
        return ".xlsx"
    if a & PPT:
        return ".pptx"
    return ".pdf"                     # tools accepting "*"


def good_sample(tool, tmp, i):
    e = ext_for(tool)
    return _MAKERS[e](tmp / f"{tool.id}_good{i}{e}")


def bad_samples(tool, tmp) -> dict:
    """Files with an ACCEPTED extension but invalid content: they get into the
    queue, then must fail cleanly during processing."""
    e = ext_for(tool)
    out = {}
    corrupt = tmp / f"{tool.id}_corrupt{e}"
    corrupt.write_bytes(os.urandom(4096))
    out["corrupt"] = corrupt
    empty = tmp / f"{tool.id}_empty{e}"
    empty.write_bytes(b"")
    out["empty"] = empty
    trunc = tmp / f"{tool.id}_truncated{e}"
    data = good_sample(tool, tmp, 99).read_bytes()
    trunc.write_bytes(data[:max(16, len(data) // 3)])
    out["truncated"] = trunc
    return out


def main() -> int:
    app = QApplication.instance() or QApplication([])
    # Modal dialogs would block a headless run; results are asserted from the
    # batch summary instead.
    QMessageBox.information = staticmethod(lambda *a, **k: None)
    QMessageBox.warning = staticmethod(lambda *a, **k: None)
    QMessageBox.question = staticmethod(lambda *a, **k: QMessageBox.No)

    saved_out = settings.output_dir
    from mico360.core.deps import find_libreoffice
    have_lo = bool(find_libreoffice())
    engine_dependent = {"office_to_pdf", "to_markdown"}
    # Content-agnostic tools work on the FILE, not its contents (File Properties
    # sets timestamps/owner), so a corrupt payload is not an error for them.
    content_agnostic = {"file_properties"}

    overrides = {
        "pdf_delete": {"pages": "1"},
        "pdf_extract": {"pages": "1"},
        "pdf_protect": {"operation": "protect", "password": "secret",
                        "confirm_password": "secret"},
        "pdf_organize": {"operation": "delete", "del_pages": "1"},
        "pdf_ocr": {"quality": "fast"},
        "pdf_convert": {"target": "image"},
    }

    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)
        out = tmp / "out"
        out.mkdir()
        settings.output_dir = str(out)
        overrides["pdf_sign"] = {"image_path": str(make_png(tmp / "_sig.png"))}

        # ============ 1. import stability ================================
        from mico360.ui.file_collector import collect_files_detailed

        big = tmp / "many"
        big.mkdir()
        for i in range(120):
            (big / f"f{i}.pdf").write_bytes(b"%PDF-1.4\n%%EOF")
        (big / "note.txt").write_text("x")
        (big / "movie.mp4").write_bytes(b"\x00" * 64)
        files, stats = collect_files_detailed([str(big)], PDF)
        check("folder import finds every supported file", len(files) == 120,
              str(len(files)))
        check("unsupported files are counted, not fatal", stats.unsupported == 2,
              str(stats.unsupported))

        files, stats = collect_files_detailed([str(big)], PDF, limit=10)
        check("import caps runaway folders instead of hanging",
              len(files) == 10 and stats.truncated,
              f"{len(files)} truncated={stats.truncated}")

        files, stats = collect_files_detailed(
            [str(tmp / "nope"), "bad:pa*th|?", str(big / "f0.pdf")], PDF)
        check("missing / invalid paths are skipped without raising",
              len(files) == 1 and stats.unreadable >= 1,
              f"{len(files)} files, unreadable={stats.unreadable}")

        deep = tmp / "deep"
        (deep / "a" / "b" / "c").mkdir(parents=True)
        (deep / "a" / "b" / "c" / "x.pdf").write_bytes(b"%PDF-1.4\n%%EOF")
        check("recursive import reaches nested folders",
              len(collect_files_detailed([str(deep)], PDF)[0]) == 1)

        # ============ 2. every tool vs bad files =========================
        for tool in TOOLS:
            if tool.id in engine_dependent and not have_lo:
                print(f"[SKIP] {tool.id}: needs LibreOffice")
                continue

            bad = bad_samples(tool, tmp)
            inputs = [str(good_sample(tool, tmp, 0))]
            if tool.mode == AGGREGATE:
                inputs.append(str(good_sample(tool, tmp, 1)))
            inputs += [str(p) for p in bad.values()]

            page = ToolPage(tool)
            page.add_paths(inputs)
            for k, v in overrides.get(tool.id, {}).items():
                _set(page, k, v)

            n_queued = len(page.items)
            s = run(page, out)

            if s.get("_nostart"):
                check(f"{tool.id}: mixed good/bad batch runs", False, "did not start")
                continue
            check(f"{tool.id}: batch with corrupt files completes (no hang)",
                  not s.get("_timeout"),
                  "TIMED OUT" if s.get("_timeout") else "")
            if s.get("_timeout"):
                continue

            done = s.get("ok", 0) + s.get("failed", 0) + s.get("skipped", 0)
            check(f"{tool.id}: every queued file is accounted for",
                  done >= 1 and s.get("total", 0) == done,
                  f"total={s.get('total')} ok={s.get('ok')} failed={s.get('failed')}")

            msgs = [e for _, e in (s.get("errors") or [])]
            if tool.id in content_agnostic:
                # Correct behaviour is to succeed on every file, including ones
                # with corrupt contents — just never crash or lose a row.
                check(f"{tool.id}: processes any file regardless of contents",
                      s.get("ok", 0) == n_queued and s.get("failed", 0) == 0,
                      f"ok={s.get('ok')} of {n_queued}")
            elif tool.mode == AGGREGATE:
                check(f"{tool.id}: aggregate run ends cleanly with a message",
                      (not s.get("failed")) or all(m.strip() for m in msgs),
                      str(msgs[:1]))
            else:
                check(f"{tool.id}: the valid file still succeeded",
                      s.get("ok", 0) >= 1, f"ok={s.get('ok')}")
                check(f"{tool.id}: bad files were flagged, not silently passed",
                      s.get("failed", 0) >= 1, f"failed={s.get('failed')}")
                check(f"{tool.id}: each failure has a readable message",
                      bool(msgs) and all(m.strip() and "Traceback" not in m
                                         for m in msgs), str(msgs[:1]))
                failed_rows = sum(1 for it in page.items if it.state == "failed")
                done_rows = sum(1 for it in page.items if it.state == "done")
                check(f"{tool.id}: queue rows show the right status",
                      failed_rows >= 1 and done_rows >= 1
                      and failed_rows + done_rows == n_queued,
                      f"{done_rows} done / {failed_rows} failed of {n_queued}")

        # ============ 3. all-bad queue must not crash ====================
        tool = next(t for t in TOOLS if t.id == "pdf_compress")
        page = ToolPage(tool)
        page.add_paths([str(p) for p in bad_samples(tool, tmp).values()])
        n = len(page.items)
        s = run(page, out)
        check("a queue of ONLY bad files finishes cleanly",
              not s.get("_timeout") and s.get("failed", 0) == n
              and s.get("ok", 0) == 0,
              f"ok={s.get('ok')} failed={s.get('failed')} of {n}")

    settings.output_dir = saved_out
    print()
    if failures:
        print(f"{len(failures)} robustness check(s) FAILED: "
              f"{', '.join(failures[:10])}")
        return 1
    print("All robustness checks passed.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
