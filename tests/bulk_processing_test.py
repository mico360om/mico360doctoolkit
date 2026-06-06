"""End-to-end BULK processing test through the real multi-threaded engine.

Drives mico360.core.engine.BatchController exactly as the app does — many files
at once, concurrent workers — and asserts the whole batch finishes with ZERO
failures. Focuses on the conversions that were failing in the field
(Office -> PDF, Document -> Markdown, incl. legacy binary .doc/.xls/.ppt) plus
lossless compression, so we know bulk runs are clean with the bundled engine.

Run:  python tests/bulk_processing_test.py
"""
from __future__ import annotations

import io
import os
import sys
import tempfile
from pathlib import Path

os.environ.setdefault("QT_QPA_PLATFORM", "offscreen")
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

from PySide6.QtCore import QEventLoop, QTimer  # noqa: E402
from PySide6.QtWidgets import QApplication  # noqa: E402

from mico360.core import deps, processors as P  # noqa: E402
from mico360.core.engine import BatchController  # noqa: E402
from mico360.core.tools import TOOLS_BY_ID  # noqa: E402

failures: list[str] = []


def check(name, ok, detail=""):
    print(f"[{'PASS' if ok else 'FAIL'}] {name}" + (f": {detail}" if detail else ""))
    if not ok:
        failures.append(name)


def run_batch(tool_id, inputs, out_dir, options=None, timeout_ms=180_000):
    """Run a real batch and return its summary dict (blocks on a Qt loop)."""
    tool = TOOLS_BY_ID[tool_id]
    ctrl = BatchController()
    loop = QEventLoop()
    box = {}

    def done(summary):
        box["summary"] = summary
        loop.quit()

    ctrl.finished.connect(done)
    guard = QTimer()
    guard.setSingleShot(True)
    guard.timeout.connect(loop.quit)
    guard.start(timeout_ms)
    ctrl.start(tool, list(inputs), out_dir, dict(options or {}), same_as_source=False)
    loop.exec()
    return box.get("summary")


def make_docx(path, text):
    import docx
    d = docx.Document()
    d.add_heading(f"Doc {path.stem}", 0)
    d.add_paragraph(text)
    t = d.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "A"; t.cell(0, 1).text = "B"
    t.cell(1, 0).text = "1"; t.cell(1, 1).text = "2"
    d.save(str(path))


def make_xlsx(path):
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["Name", "Score"])
    ws.append(["Alice", 91]); ws.append(["Bob", 88])
    wb.save(str(path))


def make_pptx(path):
    import pptx
    pr = pptx.Presentation()
    s = pr.slides.add_slide(pr.slide_layouts[1])
    s.shapes.title.text = f"Deck {path.stem}"
    s.placeholders[1].text = "Bullet one\nBullet two"
    pr.save(str(path))


def make_pdf(path, n=2, image=True):
    import fitz
    from PIL import Image
    d = fitz.open()
    for i in range(n):
        pg = d.new_page()
        pg.insert_text((72, 90), f"PDF {path.stem} page {i + 1}", fontsize=13)
        if image:
            buf = io.BytesIO()
            Image.new("RGB", (400, 300), (40, 80, 160)).save(buf, "PNG")
            pg.insert_image(fitz.Rect(72, 120, 360, 340), stream=buf.getvalue())
    d.save(str(path)); d.close()


def make_png(path, c):
    from PIL import Image
    Image.new("RGB", (600, 450), c).save(path)


def main() -> int:
    app = QApplication.instance() or QApplication(sys.argv)
    soffice = deps.find_libreoffice()
    print(f"LibreOffice engine: {soffice or 'NOT FOUND'}\n")

    td = Path(tempfile.mkdtemp(prefix="mico360_bulk_"))
    out = td / "out"; out.mkdir()
    opt = {"overwrite": True}

    # ---- Build a realistic mixed bulk set -------------------------------
    office_inputs: list[Path] = []
    for i in range(6):
        p = td / f"letter_{i}.docx"; make_docx(p, f"Body text number {i}, the quick brown fox.")
        office_inputs.append(p)
    for i in range(2):
        p = td / f"sheet_{i}.xlsx"; make_xlsx(p); office_inputs.append(p)
    for i in range(2):
        p = td / f"deck_{i}.pptx"; make_pptx(p); office_inputs.append(p)

    # Legacy binary formats (the field-failure case) — produced via LibreOffice.
    legacy: list[Path] = []
    if soffice:
        base_docx = office_inputs[0]
        for ext in ("doc", "rtf", "odt"):
            try:
                produced = P._lo_convert_to(soffice, base_docx, ext)
                dest = td / f"legacy_{ext}.{ext}"
                produced.replace(dest)
                legacy.append(dest)
            except Exception as exc:  # noqa: BLE001
                print(f"  (could not synth .{ext}: {exc})")
    office_inputs += legacy

    # ---- 1. BULK Office -> PDF (the previously-failing path) -------------
    if soffice:
        s = run_batch("office_to_pdf", office_inputs, out, opt)
        check(f"BULK Office→PDF: {len(office_inputs)} files, 0 failures",
              s and s["failed"] == 0 and s["ok"] == len(office_inputs),
              s and f"ok={s['ok']} failed={s['failed']} skipped={s['skipped']} "
                    f"errs={s['errors'][:2]}")
        check("BULK Office→PDF: every output PDF exists & non-empty",
              s and all(o.exists() and o.stat().st_size > 0 for o in s["outputs"]),
              s and f"{len(s['outputs'])} outputs")
    else:
        check("LibreOffice present for Office→PDF bulk test", False,
              "engine not found — bundle/stage it first")

    # ---- 2. BULK Document -> Markdown -----------------------------------
    md_inputs = [p for p in office_inputs if p.suffix.lower() in
                 (".docx", ".doc", ".rtf", ".odt", ".xlsx", ".pptx")]
    if soffice:
        s = run_batch("to_markdown", md_inputs, out, opt)
        check(f"BULK Document→Markdown: {len(md_inputs)} files, 0 failures",
              s and s["failed"] == 0,
              s and f"ok={s['ok']} failed={s['failed']} errs={s['errors'][:2]}")
        check("BULK Document→Markdown: every .md non-empty",
              s and all(o.exists() and o.stat().st_size > 0 for o in s["outputs"]))

    # ---- 3. BULK lossless PDF compression (24 files) --------------------
    pdfs = []
    for i in range(24):
        p = td / f"doc_{i}.pdf"; make_pdf(p, n=3, image=True); pdfs.append(p)
    s = run_batch("pdf_compress", pdfs, out, {**opt, "level": "lossless"})
    check("BULK Compress PDF (lossless): 24 files, 0 failures",
          s and s["failed"] == 0 and s["ok"] == 24,
          s and f"ok={s['ok']} failed={s['failed']} errs={s['errors'][:2]}")
    # Content must remain identical for every compressed PDF. Match each output
    # back to its source by stem (the tool adds a "_compressed" suffix), then
    # verify strict content identity.
    by_stem = {p.stem: p for p in pdfs}
    ident = checked = 0
    if s:
        for o in s["outputs"]:
            src = by_stem.get(o.stem) or by_stem.get(o.stem.replace("_compressed", ""))
            if src is None:
                continue
            checked += 1
            ok, _ = P.verify_pdf_integrity(src, o, mode="strict")
            ident += int(ok)
    check("BULK Compress PDF (lossless): all outputs verified content-identical",
          checked == len(pdfs) and ident == len(pdfs),
          f"{ident}/{checked} identical (of {len(pdfs)})")

    # ---- 4. BULK image compression (20 images) --------------------------
    imgs = []
    for i in range(20):
        p = td / f"img_{i}.png"; make_png(p, (i * 12 % 255, 70, 120)); imgs.append(p)
    s = run_batch("image_compress", imgs, out, {**opt, "level": "lossless"})
    check("BULK Compress Image (lossless): 20 files, 0 failures",
          s and s["failed"] == 0 and s["ok"] == 20,
          s and f"ok={s['ok']} failed={s['failed']} errs={s['errors'][:2]}")

    # ---- 5. Re-run Office→PDF immediately (engine reuse / no lock hang) --
    if soffice:
        s = run_batch("office_to_pdf", office_inputs[:5], out, opt)
        check("Repeated bulk run settles cleanly (no Office-lock hang)",
              s and s["failed"] == 0, s and f"failed={s['failed']}")

    import shutil
    shutil.rmtree(td, ignore_errors=True)
    print()
    if failures:
        print(f"{len(failures)} bulk check(s) FAILED: {', '.join(failures)}")
        return 1
    print("All bulk-processing checks passed — batches complete with zero failures.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
