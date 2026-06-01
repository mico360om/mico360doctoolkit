"""Tests for the v3 features: Word->PDF engine chain (fallback), compress-to-
target-size (image + PDF), and editable PDF->PowerPoint text boxes.

Run:  python tests/v3_features_test.py   (offscreen; no GUI shown)
"""
from __future__ import annotations

import os
import sys
import tempfile
from pathlib import Path

os.environ.setdefault("QT_QPA_PLATFORM", "offscreen")
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

failures: list[str] = []
rep = lambda m: None  # noqa: E731


def check(name: str, ok: bool, detail: str = "") -> None:
    print(f"[{'PASS' if ok else 'FAIL'}] {name}" + (f": {detail}" if detail else ""))
    if not ok:
        failures.append(name)


def _make_docx(path: Path) -> Path:
    import docx
    d = docx.Document()
    d.add_heading("Strategic Leadership Course", level=0)
    d.add_heading("Module 1", level=1)
    d.add_paragraph("This is a paragraph with some leadership content that should "
                    "be rendered into the PDF as selectable text.")
    p = d.add_paragraph()
    p.add_run("Bold point").bold = True
    p.add_run(" and ")
    p.add_run("italic point").italic = True
    t = d.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "Tactic"
    t.cell(0, 1).text = "Outcome"
    t.cell(1, 0).text = "Delegate"
    t.cell(1, 1).text = "Scale"
    d.save(str(path))
    return path


def _image_pdf(path: Path, photo: Path, pages: int = 3) -> Path:
    import fitz
    doc = fitz.open()
    for _ in range(pages):
        pg = doc.new_page(width=612, height=792)
        pg.insert_image(pg.rect, filename=str(photo))
    doc.save(str(path))
    doc.close()
    return path


def _text_pdf(path: Path) -> Path:
    import fitz
    doc = fitz.open()
    for i in range(2):
        pg = doc.new_page()
        pg.insert_text((72, 100), f"Slide title {i + 1}", fontsize=20)
        pg.insert_text((72, 140), "Editable bullet point one", fontsize=12)
        pg.insert_text((72, 160), "Editable bullet point two", fontsize=12)
    doc.save(str(path))
    doc.close()
    return path


def _photo(path: Path, w=1600, h=1200) -> Path:
    from PIL import Image
    img = Image.new("RGB", (w, h))
    px = img.load()
    for y in range(h):
        for x in range(w):
            px[x, y] = (x * 255 // w, y * 255 // h, (x + y) * 255 // (w + h))
    img.save(path, quality=95)
    return path


def main() -> int:
    from mico360.core import deps
    from mico360.core import processors as P

    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)
        out = tmp / "out"; out.mkdir()

        # --- Word -> PDF (engine chain) -------------------------------
        engine = ("LibreOffice" if deps.find_libreoffice()
                  else "MS Word/fallback")
        print(f"[INFO] Word->PDF will prefer: {engine}")
        docx_path = _make_docx(tmp / "Course.docx")
        res = P.word_to_pdf(docx_path, out, {}, rep)
        pdf_out = res[0]
        check("word->pdf produced a PDF", pdf_out.exists() and pdf_out.suffix == ".pdf",
              str(pdf_out))
        import fitz
        doc = fitz.open(pdf_out)
        text = "".join(pg.get_text() for pg in doc)
        doc.close()
        check("word->pdf text is selectable/searchable",
              "Strategic Leadership" in text and "Delegate" in text,
              repr(text[:60]))

        # --- image compress to target size ----------------------------
        photo = _photo(tmp / "p.jpg")
        for target_kb in (40, 120):
            r = P.image_compress(photo, out, {"level": "target", "target_kb": target_kb,
                                              "format": "keep", "overwrite": True}, rep)
            size = r[0].stat().st_size
            check(f"image target {target_kb}KB -> <= target",
                  size <= target_kb * 1024, f"{size} bytes ({size/1024:.0f} KB)")
            check(f"image target {target_kb}KB -> reasonably close (>=40%)",
                  size >= target_kb * 1024 * 0.40, f"{size/1024:.0f} KB")

        # --- PDF compress to target size ------------------------------
        ipdf = _image_pdf(tmp / "img.pdf", photo)
        src_kb = ipdf.stat().st_size / 1024
        target_kb = 200
        r = P.pdf_compress(ipdf, out, {"level": "target", "target_kb": target_kb,
                                       "overwrite": True}, rep)
        size = r[0].stat().st_size
        check(f"pdf target {target_kb}KB -> <= target (src {src_kb:.0f}KB)",
              size <= target_kb * 1024, f"{size/1024:.0f} KB")

        # --- editable PDF -> PowerPoint -------------------------------
        tpdf = _text_pdf(tmp / "deck.pdf")
        r = P.pdf_to_pptx(tpdf, out, {"overwrite": True}, rep)
        pptx_out = r[0]
        from pptx import Presentation
        prs = Presentation(str(pptx_out))
        slides = list(prs.slides)
        check("pptx has one slide per page", len(slides) == 2, str(len(slides)))
        # collect editable text from text boxes
        all_text = []
        boxes = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    boxes += 1
                    all_text.append(shape.text_frame.text)
        joined = " ".join(all_text)
        check("pptx slides contain editable text boxes", boxes >= 4, f"{boxes} boxes")
        check("pptx text is the real, editable content",
              "Slide title 1" in joined and "Editable bullet point one" in joined,
              repr(joined[:60]))

        # Empty-slide regression: an IMAGE-ONLY pdf must still yield non-empty
        # slides (page-image fallback), not blanks.
        photo2 = _photo(tmp / "scan.jpg", 1000, 700)
        imgonly = tmp / "imageonly.pdf"
        import fitz as _fitz
        _d = _fitz.open()
        for _ in range(2):
            _pg = _d.new_page(width=720, height=405)
            _pg.insert_image(_pg.rect, filename=str(photo2))
        _d.save(str(imgonly)); _d.close()
        r2 = P.pdf_to_pptx(imgonly, out, {"mode": "auto", "overwrite": True}, rep)
        prs2 = Presentation(str(r2[0]))
        pics = sum(1 for sl in prs2.slides for sh in sl.shapes if sh.shape_type == 13)
        empties = sum(1 for sl in prs2.slides if len(list(sl.shapes)) == 0)
        check("image-only PDF -> no empty slides (page-image fallback)",
              empties == 0 and pics == len(list(prs2.slides)),
              f"empties={empties} pics={pics}")

    print()
    if failures:
        print(f"{len(failures)} check(s) failed: {', '.join(failures)}")
        return 1
    print("All v3 feature checks passed.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
