"""v3.2 feature tests: remember-last-options, OCR for scanned PDFs,
and bundled-Ghostscript detection.

Run:  python tests/v32_features_test.py   (offscreen; no GUI shown)
"""
from __future__ import annotations

import os
import sys
import tempfile
from pathlib import Path

os.environ.setdefault("QT_QPA_PLATFORM", "offscreen")
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

failures: list[str] = []
rep = lambda m: None  # noqa: E731


def check(name: str, ok: bool, detail: str = "") -> None:
    print(f"[{'PASS' if ok else 'FAIL'}] {name}" + (f": {detail}" if detail else ""))
    if not ok:
        failures.append(name)


def _scanned_pdf(path: Path) -> Path:
    """A PDF whose pages are images of text (no text layer) — like a scan."""
    import fitz
    src = fitz.open()
    pg = src.new_page(width=612, height=320)
    pg.insert_text((60, 90), "Strategic Leadership Course", fontsize=26)
    pg.insert_text((60, 150), "Module One Overview", fontsize=20)
    pix = pg.get_pixmap(dpi=150)
    src.close()
    img = path.with_suffix(".png")
    pix.save(str(img))
    flat = fitz.open()
    fp = flat.new_page(width=612, height=320)
    fp.insert_image(fp.rect, filename=str(img))
    flat.save(str(path))
    flat.close()
    return path


def main() -> int:
    from PySide6.QtWidgets import QApplication
    if QApplication.instance() is None:
        QApplication([])

    # --- remember last-used options -------------------------------------
    from mico360.config import settings
    from mico360.core.tools import TOOLS_BY_ID
    from mico360.ui.options_widget import OptionsWidget

    prev = settings.tool_options("image_compress")
    settings.set_tool_options("image_compress",
                              {"level": "high", "format": "webp", "max_dimension": 640})
    ow = OptionsWidget(TOOLS_BY_ID["image_compress"])
    v = ow.values()
    check("remember-options restores saved values",
          v["level"] == "high" and v["format"] == "webp" and v["max_dimension"] == 640,
          str({k: v[k] for k in ("level", "format", "max_dimension")}))
    # save() round-trips
    ow._controls["level"].setCurrentIndex(ow._controls["level"].findData("low"))
    ow.save()
    check("remember-options save() persists",
          settings.tool_options("image_compress").get("level") == "low")
    settings.set_tool_options("image_compress", prev)  # restore

    # --- Ghostscript detection (bundled vendor path) --------------------
    from mico360.core import deps
    print(f"[INFO] Ghostscript currently: {deps.find_ghostscript()}")
    check("find_ghostscript returns a path or None",
          deps.find_ghostscript() is None or Path(deps.find_ghostscript()).exists())

    # --- OCR for scanned PDFs -------------------------------------------
    from mico360.core import processors as P
    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td); out = tmp / "out"; out.mkdir()
        scanned = _scanned_pdf(tmp / "scan.pdf")

        import fitz
        d = fitz.open(scanned)
        no_text = not any(P._page_has_text(p) for p in d)
        d.close()
        check("scanned PDF has no text layer (as expected)", no_text)

        # PDF -> PowerPoint with OCR => editable text boxes with recognised text
        r = P.pdf_to_pptx(scanned, out, {"mode": "auto", "ocr": True, "overwrite": True}, rep)
        from pptx import Presentation
        prs = Presentation(str(r[0]))
        txt = " ".join(sh.text_frame.text for sl in prs.slides for sh in sl.shapes
                       if sh.has_text_frame)
        pics = sum(1 for sl in prs.slides for sh in sl.shapes if sh.shape_type == 13)
        check("OCR pptx produced editable text boxes (not an image)",
              "Strategic" in txt and pics == 0, f"pics={pics} text={txt[:50]!r}")

        # without OCR, the same scanned PDF falls back to a page image
        r2 = P.pdf_to_pptx(scanned, out, {"mode": "auto", "ocr": False, "overwrite": True}, rep)
        prs2 = Presentation(str(r2[0]))
        pics2 = sum(1 for sl in prs2.slides for sh in sl.shapes if sh.shape_type == 13)
        check("no-OCR scanned pptx falls back to image (no empty slide)", pics2 >= 1,
              f"pics={pics2}")

        # PDF -> Word with OCR => docx containing the recognised text
        rw = P.pdf_to_word(scanned, out, {"ocr": True, "overwrite": True}, rep)
        import docx
        doc = docx.Document(str(rw[0]))
        body = " ".join(p.text for p in doc.paragraphs)
        check("OCR word produced editable text", "Strategic" in body, repr(body[:50]))

    print()
    if failures:
        print(f"{len(failures)} check(s) failed: {', '.join(failures)}")
        return 1
    print("All v3.2 feature checks passed.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
