"""v5.7 consolidated converters:
  * Office → PDF      (Word / Excel / PowerPoint, auto-detected)
  * PDF → …           (Word / PowerPoint / Excel / Image, one tool)
  * Document → Markdown (Word / Excel / PowerPoint / PDF)

Run:  python tests/consolidated_convert_test.py
"""
from __future__ import annotations

import os
import sys
import tempfile
from pathlib import Path

os.environ.setdefault("QT_QPA_PLATFORM", "offscreen")
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

failures: list[str] = []
rep = lambda *_: None  # noqa: E731


def check(name, ok, detail=""):
    print(f"[{'PASS' if ok else 'FAIL'}] {name}" + (f": {detail}" if detail else ""))
    if not ok:
        failures.append(name)


def _docx(p):
    import docx
    d = docx.Document(); d.add_heading("Title", level=1)
    d.add_paragraph("Hello world."); d.save(str(p)); return p


def _xlsx(p):
    from openpyxl import Workbook
    wb = Workbook(); ws = wb.active; ws.title = "Data"
    ws.append(["Name", "Score"]); ws.append(["Alpha", 42]); ws.append(["Beta", 7])
    wb.save(str(p)); return p


def _pptx(p):
    from pptx import Presentation
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Quarterly Review"
    s.placeholders[1].text_frame.text = "First point"
    prs.save(str(p)); return p


def _text_pdf(p, text="Searchable Heading Body"):
    import fitz
    d = fitz.open(); pg = d.new_page(width=400, height=300)
    pg.insert_text((40, 80), text, fontsize=16)
    d.save(str(p)); d.close(); return p


def test_registry():
    from mico360.core.tools import TOOLS_BY_ID, OFFICE, MARKDOWN_INPUTS
    gone = ["word_to_pdf", "excel_to_pdf", "pptx_to_pdf", "pdf_to_word",
            "pdf_to_pptx", "pdf_to_excel", "pdf_to_image", "word_to_md"]
    for tid in gone:
        check(f"old tool '{tid}' removed", tid not in TOOLS_BY_ID)
    for tid in ["pdf_convert", "office_to_pdf", "to_markdown"]:
        check(f"new tool '{tid}' present", tid in TOOLS_BY_ID)
    check("office_to_pdf accepts Word+Excel+PPT",
          TOOLS_BY_ID["office_to_pdf"].accept == OFFICE and ".xlsx" in OFFICE
          and ".pptx" in OFFICE and ".docx" in OFFICE)
    check("to_markdown accepts pdf+office",
          TOOLS_BY_ID["to_markdown"].accept == MARKDOWN_INPUTS and ".pdf" in MARKDOWN_INPUTS)
    tgt = {o.key for o in TOOLS_BY_ID["pdf_convert"].options}
    check("pdf_convert has target + per-format options",
          {"target", "word_ocr", "mode", "format", "dpi"} <= tgt)


def test_to_markdown_inputs():
    from mico360.core import processors as P
    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td); out = tmp / "o"; out.mkdir()
        # Excel → Markdown table
        r = P.to_markdown(_xlsx(tmp / "s.xlsx"), out, {"overwrite": True}, rep)
        md = r[0].read_text(encoding="utf-8")
        check("xlsx → md table header+sep",
              "| Name | Score |" in md and "| --- | --- |" in md, md[:80])
        check("xlsx → md sheet heading", "## Data" in md)
        # PowerPoint → Markdown slides
        r = P.to_markdown(_pptx(tmp / "d.pptx"), out, {"overwrite": True}, rep)
        md = r[0].read_text(encoding="utf-8")
        check("pptx → md slide heading + bullet",
              "## Slide 1" in md and "- First point" in md, md[:120])
        # PDF → Markdown text
        r = P.to_markdown(_text_pdf(tmp / "t.pdf"), out, {"overwrite": True}, rep)
        md = r[0].read_text(encoding="utf-8")
        check("pdf → md keeps text", "Searchable Heading Body" in md, md[:120])


def test_office_to_pdf_autodetect():
    from mico360.core import processors as P
    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td); out = tmp / "o"; out.mkdir()
        # .docx always works (built-in python fallback, no LibreOffice needed)
        r = P.office_to_pdf(_docx(tmp / "d.docx"), out, {"overwrite": True}, rep)
        check("office_to_pdf(.docx) → a .pdf", r and r[0].suffix == ".pdf" and r[0].exists())
        # an unsupported type raises a clear error, not a crash
        bad = tmp / "x.png"; bad.write_bytes(b"\x89PNG")
        try:
            P.office_to_pdf(bad, out, {}, rep)
            check("office_to_pdf rejects non-office", False, "no error")
        except P.ProcessError:
            check("office_to_pdf rejects non-office", True)


def test_pdf_convert_routing():
    from mico360.core import processors as P
    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td); out = tmp / "o"; out.mkdir()
        src = _text_pdf(tmp / "in.pdf")
        # → image
        r = P.pdf_convert(src, out, {"target": "image", "format": "png",
                                     "overwrite": True}, rep)
        check("pdf_convert target=image → image file",
              r and r[0].suffix.lower() == ".png" and r[0].exists())
        # → excel
        r = P.pdf_convert(src, out, {"target": "excel", "overwrite": True}, rep)
        check("pdf_convert target=excel → .xlsx", r and r[0].suffix == ".xlsx")
        # unknown target → clear error
        try:
            P.pdf_convert(src, out, {"target": "nope"}, rep)
            check("pdf_convert rejects unknown target", False, "no error")
        except P.ProcessError:
            check("pdf_convert rejects unknown target", True)


def main() -> int:
    from PySide6.QtWidgets import QApplication
    if QApplication.instance() is None:
        QApplication([])
    test_registry()
    test_to_markdown_inputs()
    test_office_to_pdf_autodetect()
    test_pdf_convert_routing()
    print()
    if failures:
        print(f"{len(failures)} check(s) failed: {', '.join(failures)}")
        return 1
    print("All consolidated-converter checks passed.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
